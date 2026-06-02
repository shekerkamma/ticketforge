#!/usr/bin/env python3
"""Use-case + realization focused deck, adapting the Canva-Pro template.

Design goals (per request):
  * NO cohort framing (no team-size / "tiny cohort" / batch-density meta).
  * Focus on USE CASES and in-depth REALIZATION detail.
  * Name the ORGANIZATIONS delivering each use case.
  * Adapt the Canva template look: teal left panel, navy accent rule, gold
    stat boxes, navy solution-stack bar, systems/users bars, navy right strip.

Run from repo root:  python3 scripts/build_yc_usecase_deck.py
"""
from __future__ import annotations

import json
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ── Canva-template palette ───────────────────────────────────────────────────
TEAL = RGBColor(0x00, 0xC9, 0xA7)
NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_2 = RGBColor(0x12, 0x24, 0x3A)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_TEAL = RGBColor(0xE0, 0xF7, 0xF1)
ACCENT = RGBColor(0x00, 0x9B, 0x82)
DARK_TEAL = RGBColor(0x00, 0x8F, 0x75)
SOFT = RGBColor(0xF4, 0xF7, 0xF8)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
GRID = RGBColor(0xD9, 0xDF, 0xE5)
GOLD = RGBColor(0xFF, 0xB8, 0x00)

HX_TEAL, HX_TEALD, HX_NAVY, HX_INK, HX_MUTED, HX_GRID, HX_GOLD = (
    "#00C9A7", "#009B82", "#0A1628", "#1B2B3C", "#5B6B7C", "#D9DFE5", "#FFB800")

FONT = "Calibri"
FONT_H = "Calibri"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.6)
CONTENT_W = SLIDE_W - MARGIN * 2
PANEL_W = Inches(8.83)
STRIP_W = SLIDE_W - PANEL_W

ANALYSIS = Path("analytics-comms/yc-agent-companies-spring-2025/analysis.json")
OUT = Path("docs/reports/yc-agent-usecases-realization.pptx")
CHARTS = Path("docs/reports/_chart_assets")
BLANK = 6
FOOTER = "Agentic AI · Business Use Cases & Realization  |  June 2026"

# ── helpers ──────────────────────────────────────────────────────────────────


def wipe(prs):
    for i in range(len(prs.slides) - 1, -1, -1):
        rid = prs.slides._sldIdLst[i].rId
        prs.part.drop_rel(rid)
        del prs.slides._sldIdLst[i]


def slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[BLANK])


def rect(s, l, t, w, h, fill, *, line=None, radius=0.0, shadow=False):
    st = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    sh = s.shapes.add_shape(st, l, t, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(1)
    if radius:
        try:
            sh.adjustments[0] = radius
        except Exception:
            pass
    sh.shadow.inherit = False
    if shadow:
        from pptx.oxml.ns import qn
        spPr = sh._element.spPr
        eff = spPr.find(qn("a:effectLst"))  # reuse the one created by inherit=False
        if eff is None:
            eff = spPr.makeelement(qn("a:effectLst"), {})
            spPr.append(eff)
        sdw = eff.makeelement(qn("a:outerShdw"), {"blurRad": "80000", "dist": "30000", "dir": "5400000", "rotWithShape": "0"})
        c = sdw.makeelement(qn("a:srgbClr"), {"val": "0A1628"})
        a = c.makeelement(qn("a:alpha"), {"val": "20000"})
        c.append(a); sdw.append(c); eff.append(sdw)
    return sh


def text(s, runs, l, t, w, h, *, size=14, color=INK, bold=False, italic=False,
         align=PP_ALIGN.LEFT, font=FONT, anchor=MSO_ANCHOR.TOP, ls=1.05,
         shrink=False, wrap=True):
    box = s.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE if shrink else MSO_AUTO_SIZE.NONE
    paras = runs if isinstance(runs, list) else [runs]
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(spec, dict):
            content = spec.get("text", "")
            ps, pc = spec.get("size", size), spec.get("color", color)
            pb, pi = spec.get("bold", bold), spec.get("italic", italic)
            pa, pf = spec.get("align", align), spec.get("font", font)
            bullet = spec.get("bullet", False)
            p.space_before = Pt(spec.get("space_before", 0))
        else:
            content, ps, pc, pb, pi, pa, pf, bullet = spec, size, color, bold, italic, align, font, False
            p.space_before = Pt(0)
        p.alignment = pa
        try:
            p.line_spacing = ls
        except Exception:
            pass
        r = p.add_run()
        r.text = ("•  " + content) if bullet else content
        r.font.name = pf
        r.font.size = Pt(ps)
        r.font.bold = pb
        r.font.italic = pi
        r.font.color.rgb = pc
    return box


def footer(s, page, total, *, dark=False):
    col = RGBColor(0x88, 0x90, 0x98) if dark else MUTED
    text(s, FOOTER, MARGIN, Inches(7.08), Inches(8), Inches(0.3), size=9, color=col)
    text(s, f"{page} / {total}", SLIDE_W - Inches(1.3), Inches(7.06), Inches(0.7), Inches(0.3),
         size=10, color=col, align=PP_ALIGN.RIGHT, bold=True)


def head(s, title, subtitle=None):
    rect(s, 0, 0, SLIDE_W, Inches(0.16), TEAL)
    text(s, title, MARGIN, Inches(0.42), CONTENT_W, Inches(0.66), size=30, color=NAVY, bold=True, font=FONT_H, shrink=True)
    rect(s, MARGIN, Inches(1.12), Inches(1.45), Inches(0.05), TEAL)
    if subtitle:
        text(s, subtitle, MARGIN, Inches(1.26), CONTENT_W, Inches(0.4), size=14.5, color=MUTED)


# ── charts ───────────────────────────────────────────────────────────────────


def chart_orgcount(clusters, path):
    data = sorted(clusters, key=lambda c: c["count"], reverse=True)
    labels = [c["name"].title() for c in data]
    vals = [c["count"] for c in data]
    fig, ax = plt.subplots(figsize=(9.6, 5.1), dpi=200)
    fig.patch.set_alpha(0); ax.set_facecolor("none")
    y = range(len(labels))
    colors = [HX_TEAL if v >= 6 else HX_TEALD for v in vals]
    ax.barh(y, vals, color=colors, height=0.6, zorder=3)
    ax.set_yticks(list(y)); ax.set_yticklabels(labels, fontsize=12.5, color=HX_INK, fontweight="bold")
    ax.invert_yaxis(); ax.set_xlim(0, max(vals) + 1.2); ax.set_xticks([])
    for i, v in enumerate(vals):
        ax.text(v + 0.15, i, str(v), va="center", fontsize=13.5, color=HX_NAVY, fontweight="bold")
    for sp in ["top", "right", "left", "bottom"]:
        ax.spines[sp].set_visible(False)
    ax.tick_params(length=0)
    plt.tight_layout(pad=0.4)
    fig.savefig(path, transparent=True, bbox_inches="tight"); plt.close(fig)


def chart_matrix(path):
    pts = [
        ("Healthcare ops", 0.80, 0.85, HX_TEAL), ("Finance & trading", 0.72, 0.74, HX_TEAL),
        ("Security & gov.", 0.40, 0.70, HX_TEALD), ("Workflow automation", 0.62, 0.60, HX_TEAL),
        ("Sales & customer ops", 0.58, 0.52, HX_TEALD), ("Data & agent infra", 0.30, 0.55, HX_TEALD),
        ("Dev tools & testing", 0.34, 0.46, HX_TEALD), ("Analytics & reporting", 0.48, 0.40, HX_MUTED),
        ("Legal & compliance", 0.70, 0.45, HX_MUTED), ("General AI apps", 0.25, 0.22, HX_MUTED),
    ]
    fig, ax = plt.subplots(figsize=(9.4, 5.3), dpi=200)
    fig.patch.set_alpha(0); ax.set_facecolor("none")
    ax.add_patch(FancyBboxPatch((0.5, 0.5), 0.5, 0.5, boxstyle="square,pad=0", facecolor=HX_TEAL, alpha=0.10, edgecolor="none"))
    ax.axhline(0.5, color=HX_GRID, lw=1.2); ax.axvline(0.5, color=HX_GRID, lw=1.2)
    for lab, x, y, c in pts:
        ax.scatter(x, y, s=240, color=c, alpha=0.92, zorder=3, edgecolors="white", linewidths=1.5)
        ax.annotate(lab, (x, y), xytext=(0, 12), textcoords="offset points", ha="center",
                    fontsize=10.5, color=HX_INK, fontweight="bold")
    ax.set_xlim(0, 1); ax.set_ylim(0, 1); ax.set_xticks([]); ax.set_yticks([])
    ax.set_xlabel("Horizontal tooling  →  Vertical workflow ownership", fontsize=11.5, color=HX_MUTED, labelpad=10)
    ax.set_ylabel("Weak  →  Strong willingness to pay", fontsize=11.5, color=HX_MUTED, labelpad=10)
    ax.text(0.985, 0.96, "Best realization payoff", ha="right", va="top", fontsize=11, color=HX_TEALD, fontweight="bold", style="italic")
    for sp in ax.spines.values():
        sp.set_color(HX_GRID)
    plt.tight_layout(pad=0.6)
    fig.savefig(path, transparent=True, bbox_inches="tight"); plt.close(fig)


# ── slide builders ───────────────────────────────────────────────────────────


def s_title(prs, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(s, PANEL_W, 0, STRIP_W, SLIDE_H, NAVY_2)
    rect(s, PANEL_W, 0, Inches(0.06), SLIDE_H, TEAL)
    # decorative dots strip
    text(s, "◆", PANEL_W + Inches(0.5), Inches(0.6), Inches(3.5), Inches(0.5), size=22, color=TEAL)
    text(s, "AGENTIC AI", MARGIN, Inches(0.95), Inches(7.5), Inches(0.45), size=16, color=TEAL, bold=True)
    text(s, "Business Use Cases", MARGIN, Inches(1.5), Inches(8.1), Inches(1.05), size=47, color=WHITE, bold=True, font=FONT_H, shrink=True)
    text(s, "& How They Get Realized", MARGIN, Inches(2.45), Inches(7.8), Inches(0.9), size=34, color=TEAL, bold=True, font=FONT_H)
    rect(s, MARGIN, Inches(3.5), Inches(1.5), Inches(0.06), TEAL)
    text(s, [
        {"text": "Ten agent use cases, the organizations delivering them, and the realization", "size": 16, "color": WHITE, "space_before": 0},
        {"text": "pattern behind each — from challenge to system of action.", "size": 16, "color": WHITE, "space_before": 3},
    ], MARGIN, Inches(3.8), Inches(7.9), Inches(1.2), shrink=True)
    chips = ["Workflow", "Customer Ops", "Analytics", "Dev & Testing", "Infra", "Healthcare"]
    x = MARGIN
    for c in chips:
        w = Inches(0.26 + 0.092 * len(c))
        rect(s, x, Inches(5.2), w, Inches(0.4), ACCENT, radius=0.5)
        text(s, c, x, Inches(5.26), w, Inches(0.3), size=11.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w + Inches(0.14)
    # org names ticker on the navy strip
    text(s, "ORGANIZATIONS REFERENCED", PANEL_W + Inches(0.5), Inches(2.0), Inches(3.6), Inches(0.3), size=11, color=TEAL, bold=True)
    orgs = ["Anana", "stratify", "Cohesive", "BitBoard", "nao Labs", "Golf", "Capacitive",
            "Airweave", "Aegis", "Pelica", "Cotool", "MindFort", "Docket",
            "Propolis", "Galen AI", "Vesence"]
    text(s, [{"text": o, "size": 11.5, "color": LIGHT_TEAL, "bold": True, "space_before": 1.5} for o in orgs],
         PANEL_W + Inches(0.5), Inches(2.4), Inches(3.6), Inches(4.4))
    text(s, FOOTER, MARGIN, Inches(6.6), Inches(8), Inches(0.3), size=11, color=RGBColor(0x6B, 0x7C, 0x8C))


def s_agenda(prs, page, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    head(s, "What This Deck Covers")
    items = [
        ("01", "The use-case landscape", "Ten agent use cases and the organizations building each one"),
        ("02", "Use cases in depth", "Challenge, realization, how-it-works, stack, and named organizations"),
        ("03", "Strategic realization", "Healthcare ops, implementation automation, risks, and a 90-day plan"),
    ]
    top = Inches(1.95)
    for n, t, d in items:
        rect(s, MARGIN, top, Inches(0.72), Inches(0.72), TEAL, radius=0.16)
        text(s, n, MARGIN, top + Inches(0.12), Inches(0.72), Inches(0.5), size=22, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(s, t, Inches(1.55), top + Inches(0.02), Inches(9.5), Inches(0.45), size=22, color=NAVY, bold=True)
        text(s, d, Inches(1.55), top + Inches(0.48), Inches(10.7), Inches(0.35), size=14, color=MUTED)
        top += Inches(1.25)
    footer(s, page, total)


def s_divider(prs, page, total, num, title, sub):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(s, 0, 0, Inches(0.18), SLIDE_H, TEAL)
    text(s, f"SECTION {num}", MARGIN, Inches(2.2), Inches(4), Inches(0.4), size=14, color=TEAL, bold=True)
    text(s, title, MARGIN, Inches(2.75), Inches(11.5), Inches(1.2), size=42, color=WHITE, bold=True, font=FONT_H)
    rect(s, MARGIN, Inches(4.05), Inches(1.5), Inches(0.06), TEAL)
    text(s, sub, MARGIN, Inches(4.35), Inches(10.5), Inches(1.0), size=18, color=LIGHT_TEAL)
    footer(s, page, total, dark=True)


def s_landscape(prs, page, total, clusters):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, SOFT)
    head(s, "The Use-Case Landscape", "Ten agent use cases — each shown with the organizations delivering it")
    order = ["workflow automation", "sales and customer ops", "analytics and reporting",
             "developer tools and testing", "data and agent infrastructure", "finance and trading",
             "security and governance", "healthcare operations", "legal and compliance workflows", "general AI applications"]
    by = {c["name"]: c for c in clusters}
    cols, gap = 2, Inches(0.3)
    cw = (CONTENT_W - gap) / cols
    ch = Inches(0.95)
    top0 = Inches(1.75)
    for i, name in enumerate(order):
        c = by[name]
        r, col = divmod(i, cols)
        x = MARGIN + col * (cw + gap)
        y = top0 + r * (ch + Inches(0.12))
        rect(s, x, y, cw, ch, WHITE, line=GRID, radius=0.07, shadow=True)
        rect(s, x, y, Inches(0.1), ch, TEAL)
        text(s, c["name"].title(), x + Inches(0.28), y + Inches(0.12), cw - Inches(1.1), Inches(0.35), size=14.5, color=NAVY, bold=True, shrink=True)
        orgs = ", ".join(co["name"] for co in c.get("companies", [])[:4])
        text(s, orgs, x + Inches(0.28), y + Inches(0.5), cw - Inches(1.1), Inches(0.4), size=11, color=MUTED, shrink=True)
        rect(s, x + cw - Inches(0.72), y + Inches(0.22), Inches(0.5), Inches(0.5), NAVY, radius=0.18)
        text(s, str(c["count"]), x + cw - Inches(0.72), y + Inches(0.3), Inches(0.5), Inches(0.35), size=17, color=GOLD, bold=True, align=PP_ALIGN.CENTER)
    footer(s, page, total)


def s_chart(prs, page, total, title, subtitle, img, takeaway, *, img_w=Inches(9.4), top=Inches(1.6)):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    head(s, title, subtitle)
    pic = s.shapes.add_picture(str(img), 0, 0, width=img_w)
    pic.left = int((SLIDE_W - pic.width) / 2)
    pic.top = top
    if pic.top + pic.height > Inches(6.1):
        sc = (Inches(6.1) - top) / pic.height
        pic.width = int(pic.width * sc); pic.height = int(pic.height * sc)
        pic.left = int((SLIDE_W - pic.width) / 2)
    text(s, takeaway, MARGIN, Inches(6.35), CONTENT_W, Inches(0.55), size=15.5, color=NAVY, bold=True)
    footer(s, page, total)


def s_pipeline(prs, page, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    head(s, "The Common Realization Pattern", "Almost every use case is the same agent loop wrapped around a different job")
    steps = [
        ("01  TRIGGER", "A queue item, message, ticket, document, alert, or schedule kicks off the run."),
        ("02  CONTEXT", "Retrieval + permissioned access pull the data, history, and policy the agent needs."),
        ("03  REASON & PLAN", "The agent decomposes the job, chooses tools, and drafts the action sequence."),
        ("04  ACT VIA TOOLS", "It executes in the systems of record — CRM, EHR, repo, ticketing, ledgers."),
        ("05  GUARDRAILS", "Confidence thresholds, policy checks, and human-in-the-loop on exceptions."),
        ("06  SYSTEM OF ACTION", "The work is completed and logged — not just summarized — with an audit trail."),
    ]
    cols, gap = 3, Inches(0.3)
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    ch = Inches(1.75)
    top0 = Inches(1.85)
    for i, (t, d) in enumerate(steps):
        r, col = divmod(i, cols)
        x = MARGIN + col * (cw + gap)
        y = top0 + r * (ch + Inches(0.28))
        rect(s, x, y, cw, ch, SOFT, line=GRID, radius=0.06)
        rect(s, x, y, cw, Inches(0.5), NAVY, radius=0.06)
        text(s, t, x + Inches(0.22), y + Inches(0.1), cw - Inches(0.4), Inches(0.32), size=13, color=TEAL, bold=True)
        text(s, d, x + Inches(0.22), y + Inches(0.62), cw - Inches(0.44), ch - Inches(0.75), size=12.5, color=INK, shrink=True, ls=1.08)
        if col < cols - 1:
            text(s, "→", x + cw - Inches(0.02), y + Inches(0.55), gap + Inches(0.1), Inches(0.6), size=22, color=TEAL_DARK if False else ACCENT, bold=True, align=PP_ALIGN.CENTER)
    footer(s, page, total)


def s_stack(prs, page, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(s, 0, 0, SLIDE_W, Inches(0.16), TEAL)
    text(s, "The Reference Realization Stack", MARGIN, Inches(0.42), CONTENT_W, Inches(0.7), size=30, color=WHITE, bold=True, font=FONT_H)
    rect(s, MARGIN, Inches(1.12), Inches(1.45), Inches(0.05), TEAL)
    text(s, "The layers organizations actually assemble to ship an agent into production", MARGIN, Inches(1.26), CONTENT_W, Inches(0.4), size=14.5, color=LIGHT_TEAL)
    layers = [
        ("EXPERIENCE & DECISION", "Inboxes, dashboards, review queues, approvals — where operators stay in control.", "Anana · BitBoard · Vesence"),
        ("ORCHESTRATION & REASONING", "Planning, tool selection, multi-step execution, evals and retries.", "stratify · Docket · Propolis"),
        ("CONTEXT & RETRIEVAL", "Permissioned access to data, memory, and policy — the layer agents fail without.", "Airweave · Kaelio · Capacitive"),
        ("ACTUATION & INTEGRATION", "Writes back into CRM, EHR, repos, ledgers and ticketing systems of record.", "Cohesive · Aegis · Cua"),
        ("GOVERNANCE & SECURITY", "Identity, scoping, audit, and guardrails on what agents may touch.", "Golf · MindFort · Cotool"),
    ]
    top = Inches(1.95)
    h = Inches(0.92)
    for i, (name, desc, orgs) in enumerate(layers):
        y = top + i * (h + Inches(0.08))
        rect(s, MARGIN, y, CONTENT_W, h, NAVY_2, line=RGBColor(0x21, 0x35, 0x4A), radius=0.05)
        rect(s, MARGIN, y, Inches(0.1), h, TEAL)
        text(s, name, MARGIN + Inches(0.3), y + Inches(0.14), Inches(3.7), Inches(0.6), size=14.5, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        text(s, desc, Inches(4.7), y, Inches(5.2), h, size=12.5, color=WHITE, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        text(s, orgs, Inches(10.1), y, Inches(2.6), h, size=11.5, color=GOLD, bold=True, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
    footer(s, page, total, dark=True)


def s_usecase(prs, page, total, uc):
    """Adapted Canva uc layout: teal panel + white section cards + navy strip."""
    s = slide(prs)
    rect(s, 0, 0, PANEL_W, SLIDE_H, TEAL)
    rect(s, PANEL_W, 0, STRIP_W, SLIDE_H, NAVY)
    rect(s, PANEL_W, 0, Inches(0.06), SLIDE_H, GOLD)

    # Title block on teal
    text(s, uc["kicker"], MARGIN, Inches(0.28), PANEL_W - Inches(1.0), Inches(0.3), size=12, color=NAVY, bold=True)
    text(s, uc["title"], MARGIN, Inches(0.55), PANEL_W - Inches(1.0), Inches(0.6), size=25, color=WHITE, bold=True, font=FONT_H, shrink=True)
    rect(s, MARGIN, Inches(1.15), Inches(1.6), Inches(0.04), NAVY)

    card_w = Inches(3.78)
    lx, rx = MARGIN, MARGIN + card_w + Inches(0.22)
    # Challenge card
    rect(s, lx, Inches(1.35), card_w, Inches(1.85), WHITE, radius=0.05, shadow=True)
    text(s, "CHALLENGE", lx + Inches(0.2), Inches(1.46), card_w - Inches(0.4), Inches(0.28), size=11.5, color=ACCENT, bold=True)
    text(s, [{"text": b, "size": 11, "color": INK, "bullet": True, "space_before": 4} for b in uc["challenge"]],
         lx + Inches(0.2), Inches(1.76), card_w - Inches(0.4), Inches(1.38), shrink=True, ls=1.02)
    # Solution card
    rect(s, rx, Inches(1.35), card_w, Inches(1.85), WHITE, radius=0.05, shadow=True)
    text(s, "HOW IT'S REALIZED", rx + Inches(0.2), Inches(1.46), card_w - Inches(0.4), Inches(0.28), size=11.5, color=ACCENT, bold=True)
    text(s, [{"text": b, "size": 11, "color": INK, "bullet": True, "space_before": 4} for b in uc["solution"]],
         rx + Inches(0.2), Inches(1.76), card_w - Inches(0.4), Inches(1.38), shrink=True, ls=1.02)

    # Stat boxes
    sb_y = Inches(3.4)
    n = len(uc["stats"])
    sw = (PANEL_W - MARGIN * 2 - Inches(0.2) * (n - 1)) / n
    for i, (num, lab) in enumerate(uc["stats"]):
        x = MARGIN + i * (sw + Inches(0.2))
        rect(s, x, sb_y, sw, Inches(0.85), NAVY, radius=0.08)
        text(s, num, x, sb_y + Inches(0.08), sw, Inches(0.42), size=19, color=GOLD, bold=True, align=PP_ALIGN.CENTER, shrink=True)
        text(s, lab, x + Inches(0.08), sb_y + Inches(0.5), sw - Inches(0.16), Inches(0.3), size=9, color=WHITE, align=PP_ALIGN.CENTER, shrink=True)

    # How it works
    hiw_y = Inches(4.45)
    rect(s, MARGIN, hiw_y, PANEL_W - MARGIN * 2, Inches(1.05), WHITE, radius=0.05, shadow=True)
    text(s, "HOW IT WORKS", MARGIN + Inches(0.2), hiw_y + Inches(0.1), Inches(4), Inches(0.26), size=11.5, color=ACCENT, bold=True)
    text(s, [{"text": f"→  {step}", "size": 10.5, "color": INK, "space_before": 2} for step in uc["how"]],
         MARGIN + Inches(0.2), hiw_y + Inches(0.38), PANEL_W - MARGIN * 2 - Inches(0.4), Inches(0.62), shrink=True, ls=1.0)

    # Solution stack bar
    st_y = Inches(5.65)
    rect(s, MARGIN, st_y, PANEL_W - MARGIN * 2, Inches(0.92), NAVY, radius=0.05)
    text(s, "SOLUTION STACK", MARGIN + Inches(0.2), st_y + Inches(0.07), Inches(3), Inches(0.22), size=10.5, color=TEAL, bold=True)
    sc_w = (PANEL_W - MARGIN * 2 - Inches(0.4)) / 4
    for i, (layer, detail) in enumerate(uc["stack"][:4]):
        cx = MARGIN + Inches(0.2) + i * sc_w
        text(s, layer, cx, st_y + Inches(0.32), sc_w - Inches(0.1), Inches(0.2), size=9.5, color=TEAL, bold=True, shrink=True)
        text(s, detail, cx, st_y + Inches(0.52), sc_w - Inches(0.1), Inches(0.36), size=8.5, color=WHITE, shrink=True, ls=1.0)

    # Systems / Users bottom bars
    by_ = Inches(6.7)
    bw = (PANEL_W - MARGIN * 2 - Inches(0.2)) / 2
    rect(s, MARGIN, by_, bw, Inches(0.34), ACCENT, radius=0.04)
    text(s, "SYSTEMS:  " + "  ·  ".join(uc["systems"][:4]), MARGIN + Inches(0.15), by_, bw - Inches(0.3), Inches(0.34), size=8.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
    rect(s, MARGIN + bw + Inches(0.2), by_, bw, Inches(0.34), DARK_TEAL, radius=0.04)
    text(s, "USERS:  " + uc["users"], MARGIN + bw + Inches(0.35), by_, bw - Inches(0.3), Inches(0.34), size=8.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE, shrink=True)

    # Right navy strip: organizations
    sx = PANEL_W + Inches(0.3)
    text(s, "ORGANIZATIONS", sx, Inches(0.45), STRIP_W - Inches(0.5), Inches(0.3), size=13, color=TEAL, bold=True)
    text(s, "DELIVERING THIS", sx, Inches(0.72), STRIP_W - Inches(0.5), Inches(0.3), size=13, color=TEAL, bold=True)
    rect(s, sx, Inches(1.08), Inches(1.0), Inches(0.04), GOLD)
    oy = Inches(1.3)
    for name, desc in uc["orgs"]:
        text(s, name, sx, oy, STRIP_W - Inches(0.5), Inches(0.3), size=14.5, color=WHITE, bold=True)
        text(s, desc, sx, oy + Inches(0.3), STRIP_W - Inches(0.55), Inches(0.62), size=10.5, color=LIGHT_TEAL, shrink=True, ls=1.0)
        oy += Inches(1.02)
    footer(s, page, total, dark=True)


def s_org_spotlight(prs, page, total, title, subtitle, orgs):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, SOFT)
    head(s, title, subtitle)
    cols, gap = 3, Inches(0.28)
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    ch = Inches(1.62)
    top0 = Inches(1.8)
    for i, (name, sub, desc) in enumerate(orgs[:6]):
        r, col = divmod(i, cols)
        x = MARGIN + col * (cw + gap)
        y = top0 + r * (ch + Inches(0.2))
        rect(s, x, y, cw, ch, WHITE, line=GRID, radius=0.07, shadow=True)
        rect(s, x, y, cw, Inches(0.1), TEAL)
        text(s, name, x + Inches(0.25), y + Inches(0.22), cw - Inches(0.5), Inches(0.35), size=17, color=NAVY, bold=True, shrink=True)
        text(s, sub, x + Inches(0.25), y + Inches(0.6), cw - Inches(0.5), Inches(0.28), size=10.5, color=ACCENT, bold=True, shrink=True)
        text(s, desc, x + Inches(0.25), y + Inches(0.9), cw - Inches(0.5), Inches(0.62), size=11, color=INK, shrink=True, ls=1.04)
        top0_done = True
    footer(s, page, total)


def s_strategic(prs, page, total, title, matches, realization, controls, question):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    text(s, title, MARGIN, Inches(0.35), Inches(12), Inches(0.62), size=27, color=NAVY, bold=True, font=FONT_H, shrink=True)
    rect(s, MARGIN, Inches(1.02), Inches(1.45), Inches(0.05), NAVY)
    col_w = Inches(5.95)
    lx, rx = MARGIN, SLIDE_W - MARGIN - col_w
    # left: organizations realizing it
    rect(s, lx, Inches(1.35), col_w, Inches(3.75), WHITE, radius=0.05, shadow=True)
    text(s, "ORGANIZATIONS REALIZING THIS", lx + Inches(0.3), Inches(1.5), col_w - Inches(0.6), Inches(0.3), size=13, color=ACCENT, bold=True)
    paras = []
    for m in matches[:5]:
        paras.append({"text": m["name"], "size": 14, "color": NAVY, "bold": True, "space_before": 7})
        paras.append({"text": m.get("proof_line") or m.get("one_liner", ""), "size": 11, "color": INK, "space_before": 1})
    text(s, paras, lx + Inches(0.3), Inches(1.9), col_w - Inches(0.6), Inches(3.1), shrink=True, ls=1.02)
    # right: realization read + controls
    rect(s, rx, Inches(1.35), col_w, Inches(3.75), NAVY, radius=0.05, shadow=True)
    text(s, "REALIZATION READ", rx + Inches(0.3), Inches(1.5), col_w - Inches(0.6), Inches(0.3), size=13, color=TEAL, bold=True)
    text(s, [{"text": l, "size": 12.5, "color": WHITE, "bullet": True, "space_before": 7} for l in realization],
         rx + Inches(0.3), Inches(1.9), col_w - Inches(0.6), Inches(1.7), shrink=True, ls=1.04)
    text(s, "WHAT TO OWN", rx + Inches(0.3), Inches(3.75), col_w - Inches(0.6), Inches(0.3), size=13, color=TEAL, bold=True)
    text(s, [{"text": l, "size": 12.5, "color": LIGHT_TEAL, "bullet": True, "space_before": 6} for l in controls],
         rx + Inches(0.3), Inches(4.1), col_w - Inches(0.6), Inches(0.95), shrink=True, ls=1.04)
    rect(s, MARGIN, Inches(5.35), CONTENT_W, Inches(1.0), WHITE, radius=0.06, shadow=True)
    text(s, question, MARGIN + Inches(0.5), Inches(5.35), CONTENT_W - Inches(1.0), Inches(1.0), size=16, color=NAVY, bold=True, italic=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(s, page, total)


def s_risks(prs, page, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    head(s, "Realization Risks & Controls", "What breaks agent use cases in production — and the control that contains each")
    rows = [
        ("Context failure", "Agents act on stale, missing, or over-broad data.", "Permissioned retrieval layer; scope reads to the task."),
        ("Silent wrong actions", "An agent writes a bad record with full confidence.", "Confidence thresholds + human-in-the-loop on exceptions."),
        ("Unbounded access", "Over-privileged agents touch systems they shouldn't.", "Identity, least-privilege scoping, and full audit trails."),
        ("No measurable outcome", "Cycle-time / error-rate impact is never proven.", "Instrument baselines before rollout; report deltas."),
        ("Workflow not owned", "The agent summarizes but doesn't complete the job.", "Close the loop into the system of record end-to-end."),
    ]
    top = Inches(1.75)
    rect(s, MARGIN, top, CONTENT_W, Inches(0.5), NAVY, radius=0.04)
    for label, x, w in [("Risk", MARGIN + Inches(0.2), Inches(2.6)), ("What goes wrong", MARGIN + Inches(3.0), Inches(5.2)), ("Control", MARGIN + Inches(8.4), Inches(3.5))]:
        text(s, label, x, top, w, Inches(0.5), size=13, color=TEAL, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    y = top + Inches(0.55)
    rh = Inches(0.82)
    for i, (risk, wrong, ctrl) in enumerate(rows):
        if i % 2 == 0:
            rect(s, MARGIN, y, CONTENT_W, rh, SOFT)
        text(s, risk, MARGIN + Inches(0.2), y, Inches(2.6), rh, size=13.5, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        text(s, wrong, MARGIN + Inches(3.0), y, Inches(5.1), rh, size=12, color=INK, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        text(s, ctrl, MARGIN + Inches(8.4), y, Inches(3.5), rh, size=12, color=ACCENT, bold=True, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        y += rh
    footer(s, page, total)


def s_roadmap(prs, page, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    head(s, "A 90-Day Realization Plan", "From a chosen use case to a proven system of action")
    phases = [
        ("Pick & scope", "Days 0-30", ["Choose one workflow-heavy use case to own end-to-end.",
                                         "Map the trigger → action loop and the systems of record.",
                                         "Instrument cycle-time and error-rate baselines."]),
        ("Build the loop", "Days 30-60", ["Ship the context + retrieval layer first.",
                                            "Wire actuation into one system of record.",
                                            "Add guardrails and a human-in-the-loop queue."]),
        ("Prove ownership", "Days 60-90", ["Run against the baseline with 2-3 design partners.",
                                             "Show measurable delta vs. the manual workflow.",
                                             "Decide: deepen the wedge or expand adjacent."]),
    ]
    n, gap = 3, Inches(0.3)
    cw = (CONTENT_W - gap * (n - 1)) / n
    top, h = Inches(2.0), Inches(3.9)
    for i, (ph, when, items) in enumerate(phases):
        x = MARGIN + i * (cw + gap)
        rect(s, x, top, cw, h, SOFT, line=GRID, radius=0.05)
        rect(s, x, top, cw, Inches(0.95), NAVY, radius=0.05)
        text(s, when, x + Inches(0.3), top + Inches(0.16), cw - Inches(0.6), Inches(0.3), size=12, color=TEAL, bold=True)
        text(s, ph, x + Inches(0.3), top + Inches(0.46), cw - Inches(0.6), Inches(0.45), size=18, color=WHITE, bold=True, shrink=True)
        text(s, [{"text": it, "size": 13, "color": INK, "bullet": True, "space_before": 12} for it in items],
             x + Inches(0.32), top + Inches(1.2), cw - Inches(0.64), h - Inches(1.45), shrink=True)
        if i < n - 1:
            text(s, "→", x + cw - Inches(0.02), top + Inches(1.6), gap + Inches(0.1), Inches(0.5), size=24, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
    footer(s, page, total)


def s_reco(prs, page, total, recs):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, SOFT)
    head(s, "Recommendations", "Where to focus realization effort first")
    gap = Inches(0.3)
    cw = (CONTENT_W - gap) / 2
    ch = Inches(1.95)
    pos = [(MARGIN, Inches(1.9)), (MARGIN + cw + gap, Inches(1.9)),
           (MARGIN, Inches(1.9) + ch + Inches(0.22)), (MARGIN + cw + gap, Inches(1.9) + ch + Inches(0.22))]
    for i, (rec, (x, y)) in enumerate(zip(recs, pos), 1):
        rect(s, x, y, cw, ch, WHITE, line=GRID, radius=0.06, shadow=True)
        rect(s, x + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7), TEAL, radius=0.16)
        text(s, f"{i:02d}", x + Inches(0.3), y + Inches(0.42), Inches(0.7), Inches(0.45), size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(s, rec["title"], x + Inches(1.2), y + Inches(0.32), cw - Inches(1.5), Inches(0.8), size=18, color=NAVY, bold=True, shrink=True)
        text(s, rec["why"], x + Inches(1.2), y + Inches(1.1), cw - Inches(1.5), Inches(0.7), size=13, color=INK, shrink=True)
    footer(s, page, total)


def s_closing(prs, page, total):
    s = slide(prs)
    rect(s, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(s, 0, 0, Inches(0.18), SLIDE_H, TEAL)
    text(s, "THE BOTTOM LINE", MARGIN, Inches(0.9), Inches(6), Inches(0.4), size=14, color=TEAL, bold=True)
    text(s, "Use cases win on realization — owning the workflow end to end, not on model novelty.",
         MARGIN, Inches(1.4), Inches(11.8), Inches(1.3), size=30, color=WHITE, bold=True, font=FONT_H)
    take = [
        ("Lead with workflow-heavy use cases", "The deepest realization payoff is in operational systems of action."),
        ("Own context, governance, and actuation", "Durable agents pair reasoning with data access and write-back into systems of record."),
        ("Treat healthcare operations as a wedge", "High-friction admin workflows with strong willingness to pay — Pelica, Aegis, Galen AI."),
        ("Close the loop, then measure it", "Complete the job and prove the cycle-time / error-rate delta."),
    ]
    top = Inches(3.05)
    for i, (t, d) in enumerate(take, 1):
        rect(s, MARGIN, top, Inches(0.55), Inches(0.55), TEAL, radius=0.18)
        text(s, str(i), MARGIN, top + Inches(0.06), Inches(0.55), Inches(0.4), size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(s, t, Inches(1.35), top - Inches(0.02), Inches(11), Inches(0.4), size=17, color=WHITE, bold=True)
        text(s, d, Inches(1.35), top + Inches(0.4), Inches(11), Inches(0.45), size=13, color=LIGHT_TEAL, shrink=True)
        top += Inches(0.97)
    footer(s, page, total, dark=True)


# ── use-case realization content ─────────────────────────────────────────────

USE_CASES = [
    {
        "kicker": "USE CASE 01  ·  HORIZONTAL OPS", "title": "Workflow Automation",
        "challenge": ["High-volume, rules-heavy operational queues", "Manual multi-step handoffs across tools", "Throughput capped by headcount"],
        "solution": ["Agents run the whole queue as a system of action", "Multi-step execution inside existing tools", "No new interface for the operator to learn"],
        "how": ["Trigger from a queue item → agent pulls context and policy", "Plans and executes the multi-step task across the team's tools", "Logs the completed action; routes only exceptions to a human"],
        "stats": [("8", "organizations"), ("Whole queue", "automation scope"), ("HITL", "on exceptions")],
        "stack": [("EXPERIENCE", "Operator inbox / queue"), ("ORCHESTRATION", "Multi-step planner"), ("CONTEXT", "Process + policy memory"), ("ACTUATION", "Writes to tools of record")],
        "systems": ["CRM", "Ticketing", "ERP", "Email/Slack"], "users": "COO · ops leaders · workflow owners",
        "orgs": [("Anana", "The agentic system of action for hotel groups and management companies."),
                 ("stratify", "AI agents that help you understand your customers' needs."),
                 ("Scalar Field", "Your agentic trading desk — agentic operations for finance."),
                 ("Cotool", "AI agents for security operations teams.")],
    },
    {
        "kicker": "USE CASE 02  ·  REVENUE", "title": "Sales & Customer Ops",
        "challenge": ["Repetitive customer-facing work at scale", "Slow response and uneven follow-up", "Reps stretched across intake and triage"],
        "solution": ["Agents qualify, route, support, and follow up", "Anchored in the revenue team's workflow", "Exploits the structured data conversations create"],
        "how": ["Inbound message / lead triggers the agent", "It enriches, qualifies, and routes against SLAs", "Drafts or sends the response; updates the CRM record"],
        "stats": [("7", "organizations"), ("Intake→follow-up", "owned workflow"), ("SLA", "anchored")],
        "stack": [("EXPERIENCE", "Rep console / CRM"), ("ORCHESTRATION", "Qualify + route logic"), ("CONTEXT", "Account + convo history"), ("ACTUATION", "CRM + messaging write-back")],
        "systems": ["CRM", "Helpdesk", "Email", "Calendar"], "users": "Sales · RevOps · CS · service managers",
        "orgs": [("Cohesive", "The agentic CRM for blue-collar services."),
                 ("Anana", "Revenue, sales, and operations working together with AI."),
                 ("Alpha Research", "Open-source, agentic knowledge bases for customer teams."),
                 ("Parsewise", "Multi-document processing for risk and ops teams.")],
    },
    {
        "kicker": "USE CASE 03  ·  KNOWLEDGE WORK", "title": "Analytics & Reporting",
        "challenge": ["Manual synthesis of raw data and context", "High-frequency, repetitive reporting", "Decisions wait on slow dashboards"],
        "solution": ["Agents turn data into decision-ready reporting", "Dashboards and summaries built on demand", "Paired with trusted, governed data access"],
        "how": ["Request or schedule triggers the reporting run", "Agent queries governed data and composes the view", "Delivers the dashboard / summary with sources cited"],
        "stats": [("6", "organizations"), ("On-demand", "dashboards"), ("Governed", "data access")],
        "stack": [("EXPERIENCE", "Dashboards / summaries"), ("ORCHESTRATION", "Query + compose"), ("CONTEXT", "Semantic + retrieval layer"), ("ACTUATION", "Publishes reports")],
        "systems": ["Warehouse", "BI tools", "Sheets", "APIs"], "users": "Heads of analytics · operators",
        "orgs": [("BitBoard", "Let your agent build dashboards and reporting."),
                 ("nao Labs", "Open-source analytics agent."),
                 ("Airweave", "The open-source context retrieval layer for agents."),
                 ("Alpha Research", "Agentic knowledge bases for decision support.")],
    },
    {
        "kicker": "USE CASE 04  ·  ENGINEERING", "title": "Developer Tools & Testing",
        "challenge": ["Agentic products create new failure modes", "QA can't scale linearly with headcount", "Regressions slip into releases"],
        "solution": ["Agents test, monitor, and validate systems", "Become part of the default release process", "Catch regressions before they ship"],
        "how": ["Code change / deploy triggers the test agent", "It drives the app, generates and runs cases", "Files findings into CI and the issue tracker"],
        "stats": [("6", "organizations"), ("In CI", "release loop"), ("↓ regressions", "outcome")],
        "stack": [("EXPERIENCE", "CI + PR checks"), ("ORCHESTRATION", "Test generation/exec"), ("CONTEXT", "Repo + spec access"), ("ACTUATION", "Issues + status writes")],
        "systems": ["GitHub/CI", "Browsers", "Issue tracker", "Observability"], "users": "Eng leaders · QA · platform teams",
        "orgs": [("Docket", "AI agents for web testing."),
                 ("Propolis", "Hands-off QA via intelligent browser agents."),
                 ("Golf", "Agentic AI security and governance for engineering."),
                 ("Capacitive", "The data gateway for agents in test pipelines.")],
    },
    {
        "kicker": "USE CASE 05  ·  PLATFORM", "title": "Data & Agent Infrastructure",
        "challenge": ["Agents fail without trustworthy context", "Uncontrolled system access is risky", "Every team rebuilds the same plumbing"],
        "solution": ["Provide the context + permissioning layer", "Become infrastructure, not a single app", "Horizontal wedge with platform stickiness"],
        "how": ["Agent product requests data or an action", "Infra layer enforces scope and retrieves context", "Returns permissioned context / brokered access"],
        "stats": [("5", "organizations"), ("Context layer", "the wedge"), ("Scoped", "access")],
        "stack": [("EXPERIENCE", "SDK / API surface"), ("ORCHESTRATION", "Access broker"), ("CONTEXT", "Retrieval + memory"), ("ACTUATION", "Permissioned connectors")],
        "systems": ["Data sources", "Vector stores", "IdP", "Connectors"], "users": "Platform eng · CTO staff",
        "orgs": [("Airweave", "The open-source context retrieval layer for agents."),
                 ("Kaelio", "The open-source context layer for data agents."),
                 ("Capacitive", "The data gateway for agents."),
                 ("Golf", "Agentic AI security and governance.")],
    },
    {
        "kicker": "USE CASE 06  ·  VERTICAL", "title": "Finance & Trading",
        "challenge": ["High-value workflows, direct $ consequences", "Structured rules, financially sensitive", "Errors and cycle time map straight to money"],
        "solution": ["Agents on claims, trading, and risk ops", "Easy ROI: cycle time and error rates = $", "Domain expertise + clear control boundaries"],
        "how": ["Claim / order / document triggers the agent", "It applies rules, checks data, and proposes the action", "Executes within limits; escalates edge cases"],
        "stats": [("4", "organizations"), ("$ impact", "measurable"), ("Bounded", "controls")],
        "stack": [("EXPERIENCE", "Ops review console"), ("ORCHESTRATION", "Rules + reasoning"), ("CONTEXT", "Docs + records"), ("ACTUATION", "Ledger / claims write")],
        "systems": ["Claims systems", "Market data", "Docs", "Ledgers"], "users": "Finance ops · claims · risk teams",
        "orgs": [("Aegis", "AI agents to win denied health-insurance claims."),
                 ("Avallon AI", "AI agents for insurance claims operations."),
                 ("Parsewise", "Multi-document processing for risk teams."),
                 ("Pelica", "Healthcare operations with financial consequences.")],
    },
    {
        "kicker": "USE CASE 07  ·  TRUST", "title": "Security & Governance",
        "challenge": ["Every agent deployment adds governance need", "Alerts, access, and policy at scale", "Security pain is urgent and measurable"],
        "solution": ["Automate security workflows", "Govern how agents touch systems and data", "Sit in a real enforcement / response loop"],
        "how": ["Alert / access request triggers the agent", "It investigates, correlates, and applies policy", "Enforces or escalates; records the decision"],
        "stats": [("4", "organizations"), ("Enforcement", "loop"), ("Audited", "decisions")],
        "stack": [("EXPERIENCE", "SOC / policy console"), ("ORCHESTRATION", "Investigate + decide"), ("CONTEXT", "Telemetry + policy"), ("ACTUATION", "Enforce / quarantine")],
        "systems": ["SIEM", "IdP", "Cloud APIs", "Ticketing"], "users": "CISO orgs · security eng · compliance",
        "orgs": [("Golf", "Agentic AI security and governance."),
                 ("MindFort", "Autonomous security agents."),
                 ("Cotool", "AI agents for security operations teams."),
                 ("Capacitive", "The data gateway controlling agent access.")],
    },
    {
        "kicker": "USE CASE 08  ·  VERTICAL", "title": "Healthcare Operations",
        "challenge": ["Repetitive patient & provider admin work", "Overloaded scheduling, billing, claims", "High friction, high willingness to pay"],
        "solution": ["Reduce manual coordination end to end", "Own a painful workflow, not a thin wrapper", "Measurable ROI in staffing and collections"],
        "how": ["Call / referral / claim triggers the agent", "It handles scheduling, intake, or billing steps", "Writes back to the EHR / RCM; flags exceptions"],
        "stats": [("3", "organizations"), ("End-to-end", "workflow"), ("ROI", "staffing + collections")],
        "stack": [("EXPERIENCE", "Front-desk / RCM UI"), ("ORCHESTRATION", "Scheduling + claims"), ("CONTEXT", "EHR + claims + labs"), ("ACTUATION", "EHR / RCM write-back")],
        "systems": ["EHR", "Claims/RCM", "Scheduling", "Telephony"], "users": "Clinic ops · RCM leaders · provider groups",
        "orgs": [("Pelica", "AI co-pilots for value-based care quality teams."),
                 ("Aegis", "AI agents to win denied health-insurance claims."),
                 ("Galen AI", "Personal AI healthcare agent on clinical + wearable data.")],
    },
    {
        "kicker": "USE CASE 09  ·  PROFESSIONAL", "title": "Legal & Compliance",
        "challenge": ["Document- and policy-heavy review work", "Expert time is expensive; throughput matters", "Language-heavy work, now tractable"],
        "solution": ["Embed agents inside the pro workflow", "Specialize around a precise task", "Save expert time, not generic 'legal AI'"],
        "how": ["Document / matter triggers the agent", "It drafts, reviews, or enforces against policy", "Surfaces the result inside the pro's tools"],
        "stats": [("2", "organizations"), ("In-workflow", "embedded"), ("Specialist", "wedge")],
        "stack": [("EXPERIENCE", "MS Office / matter UI"), ("ORCHESTRATION", "Draft + review"), ("CONTEXT", "Docs + policy corpus"), ("ACTUATION", "Edits / filings")],
        "systems": ["MS Office", "DMS", "Filing systems", "Email"], "users": "Legal ops · compliance · IP teams",
        "orgs": [("Vesence", "Cursor for lawyers — agentic AI in MS Office."),
                 ("Third Chair", "AI-native IP enforcement agent.")],
    },
]


def main():
    a = json.loads(ANALYSIS.read_text(encoding="utf-8"))
    clusters = a["use_case_clusters"]
    sync2 = a["adjacency"]["sync2_matches"]
    reprises = a["adjacency"]["reprisesai_matches"]
    recs = a["recommendations"]

    CHARTS.mkdir(parents=True, exist_ok=True)
    p_orgcount = CHARTS / "orgcount.png"
    p_matrix = CHARTS / "matrix_uc.png"
    chart_orgcount(clusters, p_orgcount)
    chart_matrix(p_matrix)

    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    wipe(prs)

    TOTAL = 27
    p = [0]

    def nx():
        p[0] += 1
        return p[0]

    s_title(prs, TOTAL); nx()
    s_agenda(prs, nx(), TOTAL)
    s_divider(prs, nx(), TOTAL, "01", "The Use-Case Landscape", "Ten agent use cases and the organizations building each one.")
    s_landscape(prs, nx(), TOTAL, clusters)
    s_chart(prs, nx(), TOTAL, "Where Organizations Are Building", "Number of named organizations active in each use case", p_orgcount,
            "Workflow automation, customer ops, analytics, testing, and infrastructure see the most activity.")
    s_pipeline(prs, nx(), TOTAL)
    s_stack(prs, nx(), TOTAL)
    s_divider(prs, nx(), TOTAL, "02", "Use Cases In Depth", "Each use case: the challenge, how it's realized, how it works, and who delivers it.")
    for uc in USE_CASES:
        s_usecase(prs, nx(), TOTAL, uc)
    s_org_spotlight(prs, nx(), TOTAL, "Organization Spotlights — Horizontal", "Operators and tooling building agent systems of action",
                    [("Anana", "Workflow automation", "The agentic system of action for hotel groups and management companies."),
                     ("Cohesive", "Sales & customer ops", "The agentic CRM for blue-collar services."),
                     ("BitBoard", "Analytics & reporting", "Let your agent build dashboards and reporting."),
                     ("Docket", "Dev tools & testing", "AI agents for web testing."),
                     ("Propolis", "Dev tools & testing", "Hands-off QA via intelligent browser agents."),
                     ("nao Labs", "Analytics & reporting", "Open-source analytics agent.")])
    s_org_spotlight(prs, nx(), TOTAL, "Organization Spotlights — Vertical & Infra", "Specialist and platform organizations",
                    [("Pelica", "Healthcare ops", "AI co-pilots for value-based care quality teams."),
                     ("Aegis", "Finance / healthcare", "AI agents to win denied health-insurance claims."),
                     ("Golf", "Security & infra", "Agentic AI security and governance."),
                     ("Airweave", "Agent infrastructure", "The open-source context retrieval layer for agents."),
                     ("Vesence", "Legal & compliance", "Cursor for lawyers — agentic AI in MS Office."),
                     ("Galen AI", "Healthcare ops", "Personal AI healthcare agent on clinical + wearable data.")])
    s_divider(prs, nx(), TOTAL, "03", "Strategic Realization", "Two high-value lanes, the risks to contain, and the plan to ship.")
    s_strategic(prs, nx(), TOTAL, "Healthcare Operations — Deep Realization", sync2,
                ["The lane is real; the strong names own deep workflow, not a front-desk wrapper.",
                 "Scheduling, claims, billing, intake, and patient comms are the system-of-action wedges.",
                 "Winners capture more of the actual clinic workflow end to end."],
                ["Own the EHR / RCM write-back, not just the conversation.",
                 "Instrument collections and staffing ROI from day one."],
                "Does the product own enough of the clinic workflow to become system-of-action software?")
    s_strategic(prs, nx(), TOTAL, "Implementation Automation — Realization Risk", reprises,
                ["Fewer direct 'AI agency' clones than expected.",
                 "The real risk: productized implementation software absorbing repeatable service work.",
                 "Generic AI delivery is easily copied; a narrow owned workflow is not."],
                ["Specialize the wedge around one repeatable, painful workflow.",
                 "Turn the service motion into owned product before someone else does."],
                "Can a service motion defend a wedge before software productizes the work?")
    s_chart(prs, nx(), TOTAL, "Where The Realization Payoff Concentrates", "Workflow ownership vs. willingness to pay, by use case", p_matrix,
            "Healthcare ops and finance sit top-right: deep workflow ownership and strong willingness to pay.",
            img_w=Inches(9.0))
    s_risks(prs, nx(), TOTAL)
    s_roadmap(prs, nx(), TOTAL)
    s_reco(prs, nx(), TOTAL, recs[:4])
    s_closing(prs, nx(), TOTAL)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Saved {OUT} with {len(prs.slides._sldIdLst)} slides (declared total {TOTAL})")


if __name__ == "__main__":
    main()
