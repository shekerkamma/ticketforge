#!/usr/bin/env python3
"""Build a clean, full-depth (28-slide) branded YC agent-companies deck.

Rewrite of build_yc_branded_pptx.py that fixes the text-overlap bugs by:
  * sizing every text box to its content and never stacking boxes that can
    overflow into each other,
  * enabling shrink-to-fit auto_size on body copy in constrained cards,
  * rendering charts with matplotlib (crisp, no hand-drawn overlap),
  * generous, consistent layout grid.

Run from repo root:  python3 scripts/build_yc_deck_v2.py
"""

from __future__ import annotations

import io
import json
from pathlib import Path
from typing import Any

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Emu, Inches, Pt

# ----------------------------------------------------------------------------
# Brand palette
# ----------------------------------------------------------------------------
NAVY = RGBColor(0x0A, 0x16, 0x28)
NAVY_2 = RGBColor(0x12, 0x24, 0x3A)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
TEAL_DARK = RGBColor(0x00, 0x9B, 0x82)
CREAM = RGBColor(0xEE, 0xEC, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
GRID = RGBColor(0xD9, 0xDF, 0xE5)
CORAL = RGBColor(0xF5, 0x47, 0x76)
GREEN = RGBColor(0x2D, 0xC4, 0x8D)

# matplotlib hex equivalents
HX_NAVY = "#0A1628"
HX_TEAL = "#00C9A7"
HX_TEALD = "#009B82"
HX_CREAM = "#EEECE1"
HX_INK = "#1B2B3C"
HX_MUTED = "#5B6B7C"
HX_GRID = "#D9DFE5"

FONT = "Aptos"
FONT_H = "Aptos Display"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
MARGIN = Inches(0.62)
CONTENT_W = SLIDE_W - MARGIN * 2

TEMPLATE_PATH = Path(
    "/mnt/c/Users/sheke/OneDrive/Desktop/Prasad_Agentic_AI_Use_Cases_Across_Industries.pptx"
)
ANALYSIS_PATH = Path("analytics-comms/yc-agent-companies-spring-2025/analysis.json")
OUTPUT_PATH = Path("docs/reports/yc-agent-companies-spring-2025-v2.pptx")
CHART_DIR = Path("docs/reports/_chart_assets")
BLANK_LAYOUT_INDEX = 6
FOOTER = "TicketForge  |  YC Agent Companies  |  June 2026"

# ----------------------------------------------------------------------------
# Low-level helpers
# ----------------------------------------------------------------------------


def remove_all_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[idx]


def blank(prs: Presentation):
    return prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])


def rect(slide, left, top, width, height, fill, *, line=None, radius=0.0, shadow=False):
    shape_type = (
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    )
    shp = slide.shapes.add_shape(shape_type, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(1)
    if radius:
        try:
            shp.adjustments[0] = radius
        except (IndexError, ValueError):
            pass
    shp.shadow.inherit = False
    if shadow:
        from pptx.oxml.ns import qn

        el = shp._element.spPr
        effect = el.find(qn("a:effectLst"))  # reuse the one created by inherit=False
        if effect is None:
            effect = el.makeelement(qn("a:effectLst"), {})
            el.append(effect)
        sdw = effect.makeelement(
            qn("a:outerShdw"),
            {"blurRad": "90000", "dist": "38100", "dir": "5400000", "rotWithShape": "0"},
        )
        clr = sdw.makeelement(qn("a:srgbClr"), {"val": "0A1628"})
        alpha = clr.makeelement(qn("a:alpha"), {"val": "22000"})
        clr.append(alpha)
        sdw.append(clr)
        effect.append(sdw)
    return shp


def text(
    slide,
    runs,
    left,
    top,
    width,
    height,
    *,
    size=16,
    color=INK,
    bold=False,
    italic=False,
    align=PP_ALIGN.LEFT,
    font=FONT,
    anchor=MSO_ANCHOR.TOP,
    line_spacing=1.04,
    space_after=0,
    shrink=False,
    wrap=True,
):
    """Add a textbox. `runs` is a str or a list of paragraph specs.

    A paragraph spec is a str, or a dict with keys: text, size, color, bold,
    italic, bullet, align, space_before.
    """
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    tf.margin_left = Emu(0)
    tf.margin_right = Emu(0)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    if shrink:
        tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    else:
        tf.auto_size = MSO_AUTO_SIZE.NONE

    paras = runs if isinstance(runs, list) else [runs]
    for i, spec in enumerate(paras):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if isinstance(spec, dict):
            content = spec.get("text", "")
            p_size = spec.get("size", size)
            p_color = spec.get("color", color)
            p_bold = spec.get("bold", bold)
            p_italic = spec.get("italic", italic)
            p_align = spec.get("align", align)
            p_font = spec.get("font", font)
            bullet = spec.get("bullet", False)
            p.space_before = Pt(spec.get("space_before", 0))
        else:
            content = spec
            p_size, p_color, p_bold, p_italic = size, color, bold, italic
            p_align, p_font, bullet = align, font, False
            p.space_before = Pt(0)
        p.alignment = p_align
        p.space_after = Pt(space_after)
        try:
            p.line_spacing = line_spacing
        except Exception:
            pass
        run = p.add_run()
        run.text = ("•  " + content) if bullet else content
        f = run.font
        f.name = p_font
        f.size = Pt(p_size)
        f.bold = p_bold
        f.italic = p_italic
        f.color.rgb = p_color
    return box


def footer(slide, page: int, *, dark=False):
    col = MUTED if not dark else RGBColor(0x7C, 0x8B, 0x9A)
    text(slide, FOOTER, MARGIN, Inches(7.06), Inches(7.5), Inches(0.3), size=10.5, color=col)
    text(
        slide,
        str(page),
        SLIDE_W - Inches(1.0),
        Inches(7.02),
        Inches(0.4),
        Inches(0.34),
        size=13,
        color=col,
        align=PP_ALIGN.RIGHT,
        bold=True,
    )


def header(slide, title, subtitle=None, *, band=True):
    """Standard light-slide header with teal accent rule."""
    if band:
        rect(slide, 0, 0, SLIDE_W, Inches(0.16), TEAL_DARK)
    text(slide, title, MARGIN, Inches(0.42), CONTENT_W, Inches(0.7), size=30, color=NAVY, bold=True, font=FONT_H)
    rect(slide, MARGIN, Inches(1.12), Inches(1.45), Inches(0.05), TEAL)
    if subtitle:
        text(slide, subtitle, MARGIN, Inches(1.26), CONTENT_W, Inches(0.4), size=15, color=MUTED)


# ----------------------------------------------------------------------------
# Charts (matplotlib -> PNG)
# ----------------------------------------------------------------------------


def _style_axes(ax):
    for spine in ["top", "right", "left"]:
        ax.spines[spine].set_visible(False)
    ax.spines["bottom"].set_color(HX_GRID)
    ax.tick_params(length=0)


def chart_density(clusters, path: Path):
    data = sorted(clusters, key=lambda c: c["count"], reverse=True)
    labels = [c["name"].title() for c in data]
    counts = [c["count"] for c in data]
    fig, ax = plt.subplots(figsize=(9.6, 5.2), dpi=200)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    y = range(len(labels))
    colors = [HX_TEAL if v >= 6 else HX_TEALD for v in counts]
    ax.barh(y, counts, color=colors, height=0.62, zorder=3)
    ax.set_yticks(list(y))
    ax.set_yticklabels(labels, fontsize=13, color=HX_INK, fontweight="bold")
    ax.invert_yaxis()
    ax.set_xlim(0, max(counts) + 1.2)
    ax.set_xticks([])
    for i, v in enumerate(counts):
        ax.text(v + 0.15, i, str(v), va="center", ha="left", fontsize=14, color=HX_NAVY, fontweight="bold")
    _style_axes(ax)
    ax.spines["bottom"].set_visible(False)
    plt.tight_layout(pad=0.4)
    fig.savefig(path, transparent=True, bbox_inches="tight")
    plt.close(fig)


def chart_team_hist(sizes, path: Path):
    bins = {"1-2": 0, "3-4": 0, "5-6": 0, "7-8": 0, "9-10": 0}
    for s in sizes:
        if s <= 2:
            bins["1-2"] += 1
        elif s <= 4:
            bins["3-4"] += 1
        elif s <= 6:
            bins["5-6"] += 1
        elif s <= 8:
            bins["7-8"] += 1
        else:
            bins["9-10"] += 1
    fig, ax = plt.subplots(figsize=(9.6, 4.9), dpi=200)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    x = range(len(bins))
    vals = list(bins.values())
    colors = [HX_TEAL if k == "1-2" else HX_TEALD for k in bins]
    ax.bar(x, vals, color=colors, width=0.66, zorder=3)
    ax.set_xticks(list(x))
    ax.set_xticklabels(list(bins.keys()), fontsize=14, color=HX_INK, fontweight="bold")
    ax.set_yticks([])
    ax.set_ylim(0, max(vals) + 2)
    for i, v in enumerate(vals):
        ax.text(i, v + 0.25, str(v), ha="center", fontsize=15, color=HX_NAVY, fontweight="bold")
    ax.set_xlabel("People per founding team", fontsize=12, color=HX_MUTED, labelpad=10)
    _style_axes(ax)
    plt.tight_layout(pad=0.5)
    fig.savefig(path, transparent=True, bbox_inches="tight")
    plt.close(fig)


def chart_matrix(path: Path):
    """2x2: x = horizontal tooling -> vertical workflow ownership;
    y = weak -> strong willingness to pay."""
    points = [
        ("Healthcare ops", 0.80, 0.85, HX_TEAL),
        ("Finance & trading", 0.72, 0.74, HX_TEAL),
        ("Security & gov.", 0.40, 0.70, HX_TEALD),
        ("Workflow automation", 0.62, 0.60, HX_TEAL),
        ("Sales & customer ops", 0.58, 0.52, HX_TEALD),
        ("Data & agent infra", 0.30, 0.55, HX_TEALD),
        ("Dev tools & testing", 0.34, 0.48, HX_TEALD),
        ("Analytics & reporting", 0.48, 0.40, HX_MUTED),
        ("Legal & compliance", 0.70, 0.45, HX_MUTED),
        ("General AI apps", 0.25, 0.22, HX_MUTED),
    ]
    fig, ax = plt.subplots(figsize=(9.4, 5.4), dpi=200)
    fig.patch.set_alpha(0)
    ax.set_facecolor("none")
    # quadrant shading
    ax.add_patch(FancyBboxPatch((0.5, 0.5), 0.5, 0.5, boxstyle="round,pad=0,rounding_size=0.01",
                                facecolor=HX_TEAL, alpha=0.10, edgecolor="none"))
    ax.axhline(0.5, color=HX_GRID, lw=1.2)
    ax.axvline(0.5, color=HX_GRID, lw=1.2)
    for label, x, y, c in points:
        ax.scatter(x, y, s=240, color=c, alpha=0.9, zorder=3, edgecolors="white", linewidths=1.5)
        ax.annotate(label, (x, y), xytext=(0, 12), textcoords="offset points",
                    ha="center", fontsize=10.5, color=HX_INK, fontweight="bold")
    ax.set_xlim(0, 1)
    ax.set_ylim(0, 1)
    ax.set_xticks([])
    ax.set_yticks([])
    ax.set_xlabel("Horizontal tooling  →  Vertical workflow ownership", fontsize=11.5,
                  color=HX_MUTED, labelpad=10)
    ax.set_ylabel("Weak  →  Strong willingness to pay", fontsize=11.5, color=HX_MUTED, labelpad=10)
    ax.text(0.985, 0.96, "Best risk-adjusted lanes", ha="right", va="top",
            fontsize=11, color=HX_TEALD, fontweight="bold", style="italic")
    for spine in ax.spines.values():
        spine.set_color(HX_GRID)
    plt.tight_layout(pad=0.6)
    fig.savefig(path, transparent=True, bbox_inches="tight")
    plt.close(fig)


# ----------------------------------------------------------------------------
# Slide builders
# ----------------------------------------------------------------------------


def s_title(prs, img_blob):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    if img_blob:
        slide.shapes.add_picture(io.BytesIO(img_blob), Inches(8.83), 0, Inches(4.5), SLIDE_H)
        rect(slide, Inches(8.83), 0, Inches(0.06), SLIDE_H, TEAL)
    text(slide, "Y COMBINATOR", MARGIN, Inches(0.95), Inches(7.5), Inches(0.4), size=16, color=TEAL, bold=True)
    text(slide, "Agent Companies", MARGIN, Inches(1.5), Inches(7.8), Inches(1.3), size=58, color=WHITE, bold=True, font=FONT_H)
    rect(slide, MARGIN, Inches(2.95), Inches(1.5), Inches(0.06), TEAL)
    text(
        slide,
        [
            {"text": "What Spring 2025's small-team agent startups are actually building", "size": 21, "color": WHITE, "bold": True, "space_before": 0},
            {"text": "Structured business use cases, a use-case density map, and the implications for Sync2 and ReprisesAI.", "size": 15, "color": CREAM, "space_before": 12},
        ],
        MARGIN,
        Inches(3.25),
        Inches(7.7),
        Inches(1.7),
    )
    chips = ["Workflow", "Customer Ops", "Healthcare", "Tooling", "Infra"]
    x = MARGIN
    for chip in chips:
        w = Inches(0.28 + 0.092 * len(chip))
        rect(slide, x, Inches(5.5), w, Inches(0.42), TEAL_DARK, radius=0.5)
        text(slide, chip, x, Inches(5.56), w, Inches(0.3), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += w + Inches(0.16)
    text(slide, FOOTER, MARGIN, Inches(6.85), Inches(8), Inches(0.3), size=11, color=RGBColor(0x6B, 0x7C, 0x8C))


def s_agenda(prs, page):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, "What You'll See")
    items = [
        ("01", "What the cohort actually says", "Headline findings, cohort size, and the use-case density map"),
        ("02", "Business use cases", "Ten clusters structured as jobs, buyers, and evidence"),
        ("03", "Strategic implications", "Sync2, ReprisesAI, recommendations, and a 90-day plan"),
    ]
    top = Inches(1.95)
    for num, title, detail in items:
        rect(slide, MARGIN, top, Inches(0.72), Inches(0.72), NAVY, radius=0.16)
        text(slide, num, MARGIN, top + Inches(0.12), Inches(0.72), Inches(0.5), size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        text(slide, title, Inches(1.6), top + Inches(0.02), Inches(9.5), Inches(0.45), size=22, color=NAVY, bold=True)
        text(slide, detail, Inches(1.6), top + Inches(0.48), Inches(10.5), Inches(0.35), size=14, color=MUTED)
        top += Inches(1.25)
    footer(slide, page)


def s_section(prs, page, num, title, subtitle):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, TEAL)
    text(slide, f"SECTION {num}", MARGIN, Inches(2.2), Inches(4), Inches(0.4), size=14, color=TEAL, bold=True)
    text(slide, title, MARGIN, Inches(2.75), Inches(11.5), Inches(1.2), size=44, color=WHITE, bold=True, font=FONT_H)
    rect(slide, MARGIN, Inches(4.05), Inches(1.5), Inches(0.06), TEAL)
    text(slide, subtitle, MARGIN, Inches(4.35), Inches(10.5), Inches(1.1), size=18, color=CREAM)
    footer(slide, page, dark=True)


def s_cards(prs, page, title, subtitle, cards):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    header(slide, title, subtitle)
    top = Inches(1.85)
    h = Inches(1.34)
    for i, (ctitle, body) in enumerate(cards, 1):
        rect(slide, MARGIN, top, CONTENT_W, h, WHITE, line=GRID, radius=0.06, shadow=True)
        rect(slide, MARGIN + Inches(0.28), top + Inches(0.27), Inches(0.8), Inches(0.8), NAVY, radius=0.18)
        text(slide, f"{i:02d}", MARGIN + Inches(0.28), top + Inches(0.42), Inches(0.8), Inches(0.5), size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        text(slide, ctitle, Inches(1.85), top + Inches(0.25), Inches(10.4), Inches(0.5), size=21, color=NAVY, bold=True)
        text(slide, body, Inches(1.85), top + Inches(0.72), Inches(10.4), Inches(0.55), size=15, color=INK, shrink=True)
        top += h + Inches(0.22)
    footer(slide, page)


def s_kpi(prs, page, title, subtitle, kpis, takeaway):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    text(slide, title, MARGIN, Inches(0.55), Inches(11.5), Inches(0.7), size=30, color=NAVY, bold=True, font=FONT_H)
    text(slide, subtitle, MARGIN, Inches(1.22), Inches(11.5), Inches(0.4), size=15, color=NAVY_2)
    n = len(kpis)
    gap = Inches(0.3)
    cw = (CONTENT_W - gap * (n - 1)) / n
    x = MARGIN
    for value, label, note in kpis:
        rect(slide, x, Inches(2.0), cw, Inches(2.1), NAVY, radius=0.07, shadow=True)
        text(slide, value, x, Inches(2.32), cw, Inches(0.85), size=46, color=TEAL, bold=True, align=PP_ALIGN.CENTER, font=FONT_H)
        text(slide, label, x + Inches(0.15), Inches(3.18), cw - Inches(0.3), Inches(0.45), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        text(slide, note, x + Inches(0.15), Inches(3.62), cw - Inches(0.3), Inches(0.42), size=12, color=CREAM, align=PP_ALIGN.CENTER, shrink=True)
        x += cw + gap
    rect(slide, MARGIN, Inches(4.65), CONTENT_W, Inches(1.4), CREAM, radius=0.06, shadow=True)
    text(slide, takeaway, MARGIN + Inches(0.6), Inches(4.65), CONTENT_W - Inches(1.2), Inches(1.4), size=21, color=NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def s_chart(prs, page, title, subtitle, img_path, takeaway, *, img_w=Inches(9.6), top=Inches(1.6)):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title, subtitle)
    pic = slide.shapes.add_picture(str(img_path), 0, 0, width=img_w)
    pic.left = int((SLIDE_W - pic.width) / 2)
    pic.top = top
    if pic.top + pic.height > Inches(6.1):
        scale = (Inches(6.1) - top) / pic.height
        pic.width = int(pic.width * scale)
        pic.height = int(pic.height * scale)
        pic.left = int((SLIDE_W - pic.width) / 2)
    text(slide, takeaway, MARGIN, Inches(6.35), CONTENT_W, Inches(0.55), size=16, color=NAVY, bold=True)
    footer(slide, page)


def s_comparison(prs, page, title, subtitle, left_title, right_title, pairs):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title, subtitle)
    col_w = Inches(5.85)
    lx = MARGIN
    rx = SLIDE_W - MARGIN - col_w
    rect(slide, lx, Inches(1.75), col_w, Inches(0.55), CORAL, radius=0.08)
    rect(slide, rx, Inches(1.75), col_w, Inches(0.55), GREEN, radius=0.08)
    text(slide, left_title, lx, Inches(1.84), col_w, Inches(0.38), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    text(slide, right_title, rx, Inches(1.84), col_w, Inches(0.38), size=16, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    top = Inches(2.5)
    h = Inches(0.84)
    for left_item, right_item in pairs:
        rect(slide, lx, top, col_w, h, RGBColor(0xFB, 0xF0, 0xF2), line=RGBColor(0xF5, 0xD7, 0xDF), radius=0.08)
        rect(slide, rx, top, col_w, h, RGBColor(0xEE, 0xFA, 0xF5), line=RGBColor(0xCC, 0xEE, 0xDE), radius=0.08)
        text(slide, [{"text": "✕   " + left_item, "color": INK}], lx + Inches(0.3), top, col_w - Inches(0.6), h, size=14.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        text(slide, [{"text": "✓   " + right_item, "color": INK}], rx + Inches(0.3), top, col_w - Inches(0.6), h, size=14.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        top += h + Inches(0.16)
    footer(slide, page)


def s_jobs_table(prs, page, title, subtitle, clusters):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title, subtitle)
    rows = [("Buyer", "buyer"), ("Business job", "business_job"), ("Value prop", "value_prop"), ("Why now", "why_now")]
    label_w = Inches(1.55)
    grid_left = MARGIN + label_w + Inches(0.12)
    n = len(clusters)
    gap = Inches(0.12)
    cw = (SLIDE_W - grid_left - MARGIN - gap * (n - 1)) / n
    # column headers
    x = grid_left
    for c in clusters:
        rect(slide, x, Inches(1.7), cw, Inches(0.62), NAVY, radius=0.08)
        text(slide, c["name"].title(), x + Inches(0.08), Inches(1.7), cw - Inches(0.16), Inches(0.62), size=12.5, color=WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        x += cw + gap
    row_tops = [Inches(2.45), Inches(3.18), Inches(4.0), Inches(4.95)]
    row_hts = [Inches(0.62), Inches(0.7), Inches(0.83), Inches(1.0)]
    for (label, field), rtop, rht in zip(rows, row_tops, row_hts):
        rect(slide, MARGIN, rtop, label_w, rht, TEAL if label != "Business job" else TEAL_DARK)
        text(slide, label, MARGIN + Inches(0.12), rtop, label_w - Inches(0.2), rht, size=12.5, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        x = grid_left
        for c in clusters:
            rect(slide, x, rtop, cw, rht, RGBColor(0xF3, 0xF6, 0xF8), line=GRID)
            text(slide, c[field], x + Inches(0.1), rtop + Inches(0.06), cw - Inches(0.2), rht - Inches(0.12), size=10.5, color=INK, shrink=True, line_spacing=1.0)
            x += cw + gap
    footer(slide, page)


def s_cluster_detail(prs, page, cluster, rank, total):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    rect(slide, 0, 0, SLIDE_W, Inches(1.55), NAVY)
    text(slide, f"USE CASE {rank} OF {total}", MARGIN, Inches(0.28), Inches(6), Inches(0.3), size=12, color=TEAL, bold=True)
    text(slide, cluster["name"].title(), MARGIN, Inches(0.6), Inches(8.5), Inches(0.8), size=30, color=WHITE, bold=True, font=FONT_H)
    # count badge
    rect(slide, SLIDE_W - MARGIN - Inches(1.7), Inches(0.42), Inches(1.7), Inches(0.78), TEAL, radius=0.12)
    text(slide, [
        {"text": str(cluster["count"]), "size": 30, "color": NAVY, "bold": True, "align": PP_ALIGN.CENTER},
    ], SLIDE_W - MARGIN - Inches(1.7), Inches(0.46), Inches(1.7), Inches(0.5), align=PP_ALIGN.CENTER, font=FONT_H)
    text(slide, "companies", SLIDE_W - MARGIN - Inches(1.7), Inches(0.94), Inches(1.7), Inches(0.25), size=11, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    fields = [
        ("BUYER", cluster["buyer"]),
        ("BUSINESS JOB", cluster["business_job"]),
        ("VALUE PROP", cluster["value_prop"]),
        ("WHY NOW", cluster["why_now"]),
    ]
    # left column: 2 cards, right column: 2 cards
    positions = [
        (MARGIN, Inches(1.95)),
        (MARGIN, Inches(3.55)),
        (Inches(6.95), Inches(1.95)),
        (Inches(6.95), Inches(3.55)),
    ]
    cw = Inches(5.78)
    ch = Inches(1.42)
    for (lab, val), (x, y) in zip(fields, positions):
        rect(slide, x, y, cw, ch, RGBColor(0xF5, 0xF7, 0xF9), line=GRID, radius=0.06)
        rect(slide, x, y, Inches(0.08), ch, TEAL)
        text(slide, lab, x + Inches(0.28), y + Inches(0.16), cw - Inches(0.5), Inches(0.3), size=12, color=TEAL_DARK, bold=True)
        text(slide, val, x + Inches(0.28), y + Inches(0.5), cw - Inches(0.5), Inches(0.82), size=14, color=INK, shrink=True)
    # recommendation strip
    rect(slide, MARGIN, Inches(5.2), CONTENT_W, Inches(0.92), NAVY, radius=0.06)
    text(slide, "RECOMMENDATION", MARGIN + Inches(0.3), Inches(5.32), Inches(2.6), Inches(0.3), size=11, color=TEAL, bold=True)
    text(slide, cluster["recommendation"], MARGIN + Inches(0.3), Inches(5.58), CONTENT_W - Inches(0.6), Inches(0.5), size=14, color=WHITE, shrink=True)
    # example companies (single box, two runs, to avoid label/value collision)
    examples = ", ".join(cluster.get("examples", [])[:5])
    box = slide.shapes.add_textbox(MARGIN, Inches(6.32), CONTENT_W, Inches(0.4))
    tf = box.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    para = tf.paragraphs[0]
    r1 = para.add_run(); r1.text = "Cohort examples:   "
    r1.font.name = FONT; r1.font.size = Pt(12); r1.font.bold = True; r1.font.color.rgb = MUTED
    r2 = para.add_run(); r2.text = examples
    r2.font.name = FONT; r2.font.size = Pt(12); r2.font.bold = True; r2.font.color.rgb = INK
    footer(slide, page)


def s_lanes_table(prs, page, title, subtitle, lanes):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title, subtitle)
    cols = [("Lane", Inches(2.7)), ("Count", Inches(1.0)), ("Business job", Inches(5.9)), ("Cohort evidence", Inches(2.5))]
    xs = [MARGIN]
    for _, w in cols[:-1]:
        xs.append(xs[-1] + w)
    # header row
    rect(slide, MARGIN, Inches(1.75), CONTENT_W, Inches(0.5), CREAM, radius=0.04)
    for (label, w), x in zip(cols, xs):
        text(slide, label, x + Inches(0.15), Inches(1.75), w - Inches(0.2), Inches(0.5), size=13, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    top = Inches(2.32)
    rh = Inches(0.86)
    for i, lane in enumerate(lanes):
        if i % 2 == 1:
            rect(slide, MARGIN, top, CONTENT_W, rh, RGBColor(0xF6, 0xF8, 0xFA))
        text(slide, lane["name"], xs[0] + Inches(0.15), top, cols[0][1] - Inches(0.2), rh, size=13.5, color=NAVY, bold=True, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        text(slide, str(lane["count"]), xs[1] + Inches(0.15), top, cols[1][1] - Inches(0.2), rh, size=15, color=TEAL_DARK, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, lane["business_job"], xs[2] + Inches(0.15), top, cols[2][1] - Inches(0.3), rh, size=12.5, color=INK, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        text(slide, ", ".join(lane["examples"][:3]), xs[3] + Inches(0.15), top, cols[3][1] - Inches(0.25), rh, size=12, color=MUTED, anchor=MSO_ANCHOR.MIDDLE, shrink=True)
        rect(slide, MARGIN, top + rh, CONTENT_W, Pt(0.75), GRID)
        top += rh
    text(slide, "Smaller clusters can still be strategically attractive when workflow pain and willingness to pay are high.",
         MARGIN, Inches(6.4), CONTENT_W, Inches(0.5), size=14.5, color=NAVY, bold=True)
    footer(slide, page)


def s_company_grid(prs, page, title, subtitle, companies):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    header(slide, title, subtitle)
    cols, rows = 3, 3
    gap = Inches(0.28)
    cw = (CONTENT_W - gap * (cols - 1)) / cols
    ch = Inches(1.5)
    start_top = Inches(1.85)
    for idx, comp in enumerate(companies[: cols * rows]):
        r, c = divmod(idx, cols)
        x = MARGIN + c * (cw + gap)
        y = start_top + r * (ch + Inches(0.18))
        rect(slide, x, y, cw, ch, WHITE, line=GRID, radius=0.07, shadow=True)
        rect(slide, x, y, cw, Inches(0.1), TEAL)
        text(slide, comp["name"], x + Inches(0.25), y + Inches(0.22), cw - Inches(1.0), Inches(0.4), size=16, color=NAVY, bold=True, shrink=True)
        text(slide, f"{comp['team_size']}p", x + cw - Inches(0.85), y + Inches(0.24), Inches(0.65), Inches(0.32), size=12, color=TEAL_DARK, bold=True, align=PP_ALIGN.RIGHT)
        text(slide, comp.get("proof_line") or comp.get("what_they_do", ""), x + Inches(0.25), y + Inches(0.66), cw - Inches(0.5), Inches(0.78), size=11.5, color=INK, shrink=True, line_spacing=1.02)
    footer(slide, page)


def s_case_study(prs, page, title, matches, implications, kpis, question):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    text(slide, title, MARGIN, Inches(0.4), Inches(11.5), Inches(0.7), size=30, color=NAVY, bold=True, font=FONT_H)
    rect(slide, MARGIN, Inches(1.08), Inches(1.45), Inches(0.05), NAVY)

    col_w = Inches(5.85)
    lx, rx = MARGIN, SLIDE_W - MARGIN - col_w
    # left card: matches
    rect(slide, lx, Inches(1.4), col_w, Inches(3.55), WHITE, radius=0.05, shadow=True)
    text(slide, "CLOSEST YC MATCHES", lx + Inches(0.3), Inches(1.6), col_w - Inches(0.6), Inches(0.35), size=14, color=TEAL_DARK, bold=True)
    paras = []
    for m in matches[:4]:
        paras.append({"text": m["name"], "size": 14.5, "color": NAVY, "bold": True, "space_before": 9})
        paras.append({"text": m.get("proof_line") or m.get("one_liner", ""), "size": 12, "color": INK, "space_before": 1})
    text(slide, paras, lx + Inches(0.3), Inches(2.05), col_w - Inches(0.6), Inches(2.8), shrink=True)
    # right card: implications
    rect(slide, rx, Inches(1.4), col_w, Inches(3.55), NAVY, radius=0.05, shadow=True)
    text(slide, "WHAT THIS MEANS", rx + Inches(0.3), Inches(1.6), col_w - Inches(0.6), Inches(0.35), size=14, color=TEAL, bold=True)
    text(slide, [{"text": line, "color": CREAM, "size": 14, "bullet": True, "space_before": 10} for line in implications],
         rx + Inches(0.3), Inches(2.05), col_w - Inches(0.6), Inches(2.8), shrink=True)
    # kpi chips
    n = len(kpis)
    gap = Inches(0.25)
    kw = (CONTENT_W - gap * (n - 1)) / n
    x = MARGIN
    for value, label in kpis:
        rect(slide, x, Inches(5.15), kw, Inches(0.95), NAVY, radius=0.1)
        text(slide, value, x, Inches(5.24), kw, Inches(0.5), size=24, color=TEAL, bold=True, align=PP_ALIGN.CENTER, font=FONT_H)
        text(slide, label, x + Inches(0.1), Inches(5.74), kw - Inches(0.2), Inches(0.32), size=11.5, color=WHITE, align=PP_ALIGN.CENTER, shrink=True)
        x += kw + gap
    text(slide, question, MARGIN, Inches(6.35), CONTENT_W, Inches(0.55), size=15.5, color=NAVY, bold=True, italic=True)
    footer(slide, page)


def s_two_col(prs, page, title, subtitle, left_head, left_items, right_head, right_items):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title, subtitle)
    col_w = Inches(5.85)
    lx, rx = MARGIN, SLIDE_W - MARGIN - col_w
    for x, head, items, accent in [(lx, left_head, left_items, GREEN), (rx, right_head, right_items, CORAL)]:
        rect(slide, x, Inches(1.8), col_w, Inches(4.5), RGBColor(0xF5, 0xF7, 0xF9), line=GRID, radius=0.04)
        rect(slide, x, Inches(1.8), col_w, Inches(0.6), accent, radius=0.04)
        text(slide, head, x + Inches(0.3), Inches(1.8), col_w - Inches(0.6), Inches(0.6), size=16, color=WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)
        text(slide, [{"text": it, "size": 14, "color": INK, "bullet": True, "space_before": 12} for it in items],
             x + Inches(0.35), Inches(2.65), col_w - Inches(0.7), Inches(3.4), shrink=True)
    footer(slide, page)


def s_rec_grid(prs, page, title, subtitle, recs):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    header(slide, title, subtitle)
    gap = Inches(0.3)
    cw = (CONTENT_W - gap) / 2
    ch = Inches(1.95)
    positions = [(MARGIN, Inches(1.9)), (MARGIN + cw + gap, Inches(1.9)),
                 (MARGIN, Inches(1.9) + ch + Inches(0.22)), (MARGIN + cw + gap, Inches(1.9) + ch + Inches(0.22))]
    for i, (rec, (x, y)) in enumerate(zip(recs, positions), 1):
        rect(slide, x, y, cw, ch, WHITE, line=GRID, radius=0.06, shadow=True)
        rect(slide, x + Inches(0.3), y + Inches(0.3), Inches(0.7), Inches(0.7), TEAL, radius=0.16)
        text(slide, f"{i:02d}", x + Inches(0.3), y + Inches(0.42), Inches(0.7), Inches(0.45), size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(slide, rec["title"], x + Inches(1.2), y + Inches(0.32), cw - Inches(1.5), Inches(0.8), size=18, color=NAVY, bold=True, shrink=True)
        text(slide, rec["why"], x + Inches(1.2), y + Inches(1.1), cw - Inches(1.5), Inches(0.7), size=13.5, color=INK, shrink=True)
    footer(slide, page)


def s_roadmap(prs, page, title, subtitle, phases):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title, subtitle)
    n = len(phases)
    gap = Inches(0.3)
    cw = (CONTENT_W - gap * (n - 1)) / n
    top = Inches(2.0)
    h = Inches(3.9)
    for i, (phase, when, items) in enumerate(phases):
        x = MARGIN + i * (cw + gap)
        rect(slide, x, top, cw, h, RGBColor(0xF5, 0xF7, 0xF9), line=GRID, radius=0.05)
        rect(slide, x, top, cw, Inches(0.95), NAVY, radius=0.05)
        text(slide, when, x + Inches(0.3), top + Inches(0.16), cw - Inches(0.6), Inches(0.3), size=12, color=TEAL, bold=True)
        text(slide, phase, x + Inches(0.3), top + Inches(0.46), cw - Inches(0.6), Inches(0.45), size=17, color=WHITE, bold=True, shrink=True)
        text(slide, [{"text": it, "size": 13, "color": INK, "bullet": True, "space_before": 11} for it in items],
             x + Inches(0.32), top + Inches(1.2), cw - Inches(0.64), h - Inches(1.45), shrink=True)
        if i < n - 1:
            text(slide, "→", x + cw - Inches(0.02), top + Inches(1.6), gap + Inches(0.1), Inches(0.5), size=24, color=TEAL_DARK, bold=True, align=PP_ALIGN.CENTER)
    footer(slide, page)


def s_methodology(prs, page, title, bullets, note):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    header(slide, title)
    text(slide, [{"text": b, "size": 16, "color": INK, "bullet": True, "space_before": 14} for b in bullets],
         MARGIN + Inches(0.1), Inches(1.7), CONTENT_W - Inches(0.2), Inches(3.4))
    rect(slide, MARGIN, Inches(5.35), CONTENT_W, Inches(1.1), CREAM, line=GRID, radius=0.06)
    text(slide, note, MARGIN + Inches(0.5), Inches(5.35), CONTENT_W - Inches(1.0), Inches(1.1), size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    footer(slide, page)


def s_closing(prs, page, takeaways):
    slide = blank(prs)
    rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    rect(slide, 0, 0, Inches(0.18), SLIDE_H, TEAL)
    text(slide, "THE BOTTOM LINE", MARGIN, Inches(0.85), Inches(6), Inches(0.4), size=14, color=TEAL, bold=True)
    text(slide, "Winners are defined by workflow choice and ownership — not model novelty.",
         MARGIN, Inches(1.35), Inches(11.8), Inches(1.3), size=32, color=WHITE, bold=True, font=FONT_H)
    top = Inches(3.1)
    for i, (t, d) in enumerate(takeaways, 1):
        rect(slide, MARGIN, top, Inches(0.55), Inches(0.55), TEAL, radius=0.18)
        text(slide, str(i), MARGIN, top + Inches(0.06), Inches(0.55), Inches(0.4), size=18, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        text(slide, t, Inches(1.35), top - Inches(0.02), Inches(11), Inches(0.4), size=17, color=WHITE, bold=True)
        text(slide, d, Inches(1.35), top + Inches(0.4), Inches(11), Inches(0.4), size=13, color=CREAM, shrink=True)
        top += Inches(0.95)
    footer(slide, page, dark=True)


# ----------------------------------------------------------------------------
# Assembly
# ----------------------------------------------------------------------------


def main() -> None:
    if not ANALYSIS_PATH.exists():
        raise SystemExit(f"Missing analysis pack: {ANALYSIS_PATH}")
    analysis = json.loads(ANALYSIS_PATH.read_text(encoding="utf-8"))
    clusters = analysis["use_case_clusters"]
    by_name = {c["name"]: c for c in clusters}
    cohort = analysis["cohort_summary"]
    companies = analysis["companies"]
    sync2 = analysis["adjacency"]["sync2_matches"]
    reprises = analysis["adjacency"]["reprisesai_matches"]
    recs = analysis["recommendations"]

    # charts
    CHART_DIR.mkdir(parents=True, exist_ok=True)
    p_density = CHART_DIR / "density.png"
    p_hist = CHART_DIR / "team_hist.png"
    p_matrix = CHART_DIR / "matrix.png"
    chart_density(clusters, p_density)
    chart_team_hist(cohort["team_sizes"], p_hist)
    chart_matrix(p_matrix)

    # template (for the title art only)
    img_blob = None
    if TEMPLATE_PATH.exists():
        prs = Presentation(str(TEMPLATE_PATH))
        try:
            img_blob = prs.slides[0].shapes[2].image.blob
        except Exception:
            img_blob = None
    else:
        prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    remove_all_slides(prs)

    p = 0

    def page():
        nonlocal p
        p += 1
        return p

    # 1 title
    s_title(prs, img_blob)
    page()
    # 2 agenda
    s_agenda(prs, page())
    # 3 exec kpi
    s_kpi(
        prs, page(), "The Cohort At A Glance", "A young, focused cohort with a visible use-case pattern",
        [
            ("27", "companies", "filtered Spring 2025 cohort"),
            ("3", "median team", "very early-stage teams"),
            ("13", "teams of 1-2", "nearly half the cohort"),
            ("10", "use-case lanes", "distinct business shapes"),
        ],
        "Likely winners will be defined by workflow choice and ownership, not by model novelty.",
    )
    # 4 section 1
    s_section(prs, page(), "01", "What The Cohort Actually Says",
              "The data is more operational and workflow-shaped than the generic agent narrative suggests.")
    # 5 three conclusions
    s_cards(
        prs, page(), "Three Conclusions Matter Most", "The headline reads of the cohort",
        [
            ("Workflow automation leads", "The strongest YC signal is not generic assistants. It is agents wrapped around concrete operational jobs."),
            ("Ops and tooling dominate", "Customer ops, analytics, QA/testing, context, and infra appear far more often than broad platform narratives."),
            ("Healthcare is a real vertical wedge", "Administrative healthcare workflows show enough density and pain to warrant focused product bets."),
        ],
    )
    # 6 team-size histogram
    s_chart(prs, page(), "This Cohort Is Extremely Early",
            "Distribution of founding-team size across the 27 companies", p_hist,
            "13 of 27 teams have just 1-2 people; the median team is 3. Execution focus beats headcount here.",
            img_w=Inches(9.2))
    # 7 density bar
    s_chart(prs, page(), "Use-Case Density Favors Real Jobs",
            "Number of companies per use-case cluster (companies can span lanes)", p_density,
            "The center of gravity is workflow automation, customer ops, analytics, testing, and infrastructure.",
            img_w=Inches(9.6))
    # 8 expected vs observed
    s_comparison(
        prs, page(), "What People Expect vs. What The Cohort Shows",
        "Previous expectation vs. what the YC cohort actually reveals", "Expected narrative", "Observed YC pattern",
        [
            ("General-purpose agent assistants", "Workflow-bound systems of action"),
            ("Broad agent-platform positioning", "Buyer-owned operational categories"),
            ("Abstract AI capability claims", "Testing, context, claims, reporting, ops"),
            ("Model novelty as the main story", "Clear job boundaries & measurable outcomes"),
        ],
    )
    # 9 section 2
    s_section(prs, page(), "02", "Business Use Cases",
              "Ten clusters, structured as the jobs they do, the buyers they serve, and the evidence behind them.")
    # 10 jobs table (top 4)
    top4 = [by_name["workflow automation"], by_name["sales and customer ops"],
            by_name["analytics and reporting"], by_name["developer tools and testing"]]
    s_jobs_table(prs, page(), "Top Use Cases, Structured As Business Jobs",
                 "The four highest-signal business shapes in the cohort", top4)
    # 11-15 cluster details (top 5 by count)
    detail_order = ["workflow automation", "sales and customer ops", "analytics and reporting",
                    "developer tools and testing", "data and agent infrastructure"]
    for i, name in enumerate(detail_order, 1):
        s_cluster_detail(prs, page(), by_name[name], i, len(detail_order))
    # 16 vertical lanes table
    lanes = [
        {"name": "Finance & Trading", "count": 4, "business_job": by_name["finance and trading"]["business_job"], "examples": by_name["finance and trading"]["examples"]},
        {"name": "Security & Governance", "count": 4, "business_job": by_name["security and governance"]["business_job"], "examples": by_name["security and governance"]["examples"]},
        {"name": "Healthcare Operations", "count": 3, "business_job": by_name["healthcare operations"]["business_job"], "examples": by_name["healthcare operations"]["examples"]},
        {"name": "General AI Applications", "count": 3, "business_job": by_name["general AI applications"]["business_job"], "examples": by_name["general AI applications"]["examples"]},
        {"name": "Legal & Compliance", "count": 2, "business_job": by_name["legal and compliance workflows"]["business_job"], "examples": by_name["legal and compliance workflows"]["examples"]},
    ]
    s_lanes_table(prs, page(), "Vertical & Specialist Lanes Matter Too",
                  "Smaller clusters, but high pain and high willingness to pay", lanes)
    # 17 2x2 matrix
    s_chart(prs, page(), "Where The Strong Lanes Sit",
            "Workflow ownership vs. willingness to pay — by use-case cluster", p_matrix,
            "Healthcare ops and finance sit top-right: deep workflow ownership and strong willingness to pay.",
            img_w=Inches(9.2))
    # 18 company evidence grid
    feat = ["Anana", "Golf", "Pelica", "stratify", "BitBoard", "Airweave", "Aegis", "Cotool", "MindFort"]
    cmap = {c["name"]: c for c in companies}
    grid = [cmap[n] for n in feat if n in cmap]
    if len(grid) < 9:
        grid += [c for c in companies if c not in grid][: 9 - len(grid)]
    s_company_grid(prs, page(), "Representative Companies", "A cross-section of the filtered cohort", grid)
    # 19 section 3
    s_section(prs, page(), "03", "Strategic Implications",
              "Where the YC pattern supports Sync2, where it pressures ReprisesAI, and what to build next.")
    # 20 sync2 case study
    s_case_study(
        prs, page(), "Sync2 Has Real YC Adjacency", sync2,
        [
            "The lane is real, but the stronger names own deeper workflow than a thin front-desk wrapper.",
            "Scheduling, claims, billing, intake, and patient comms are stronger system-of-action wedges.",
            "The winning product likely captures more of the actual clinic workflow.",
        ],
        [("5", "strong adjacency matches"), ("Ops", "workflow depth wins"), ("$$$", "high-friction buyer pain")],
        "Key question: does the product own enough workflow to become system-of-action software?",
    )
    # 21 reprises case study
    s_case_study(
        prs, page(), "ReprisesAI Faces Productization Risk", reprises,
        [
            "There are fewer direct AI-agency clones than expected.",
            "The bigger risk is productized implementation software absorbing repeatable service work.",
            "The safer posture is a specialized wedge, not generic AI delivery.",
        ],
        [("4", "clear comparables"), ("Risk", "services get productized"), ("Focus", "specialize the wedge")],
        "Key question: can a service offer defend a wedge before software productizes the work?",
    )
    # 22 support vs pressure
    s_two_col(
        prs, page(), "Support vs. Pressure, Side By Side",
        "How the cohort cuts for each strategy",
        "Where the pattern SUPPORTS Sync2",
        [
            "Healthcare ops is a validated, dense lane.",
            "Buyers feel acute, measurable admin pain.",
            "Deep workflow ownership is rewarded, not punished.",
            "Voice + RCM + intake are proven entry wedges.",
        ],
        "Where it PRESSURES ReprisesAI",
        [
            "Productized implementation tools are emerging fast.",
            "Generic 'AI agency' positioning is easily copied.",
            "Repeatable service work is the first to be automated.",
            "Defensibility requires a narrow, owned workflow.",
        ],
    )
    # 23 recommendations
    s_rec_grid(prs, page(), "What I Would Do Next", "Recommendations follow directly from the cohort structure", recs[:4])
    # 24 risks
    s_cards(
        prs, page(), "Risks & Watch-Outs", "What could invalidate this read",
        [
            ("Cohort is tiny (n=27)", "Patterns are directional, not statistically conclusive. Treat counts as signal, re-test next batch."),
            ("Labels are inferred", "Use-case clusters come from company text, not official YC categories — boundaries are fuzzy."),
            ("Cohort attrition is real", "Some names have already wound down (e.g. Capacitive), so survivorship will reshape the map."),
        ],
    )
    # 25 roadmap
    s_roadmap(
        prs, page(), "A 90-Day Action Plan", "Turning the read into concrete moves",
        [
            ("Validate the wedge", "Days 0-30", [
                "Pick one workflow-heavy lane to own end-to-end.",
                "Interview 10 buyers in healthcare ops or finance.",
                "Define the system-of-action boundary precisely.",
            ]),
            ("Build the proof", "Days 30-60", [
                "Ship a narrow agent that owns one painful queue.",
                "Instrument cycle-time and error-rate baselines.",
                "Secure 2-3 design partners with clear SLAs.",
            ]),
            ("Prove ownership", "Days 60-90", [
                "Embed into the default workflow / release loop.",
                "Show measurable ROI vs. manual baseline.",
                "Decide: deepen the wedge or expand adjacent.",
            ]),
        ],
    )
    # 26 methodology
    s_methodology(
        prs, page(), "Method & Reproducibility",
        [
            "Website tested: https://www.ycombinator.com/companies",
            "yc-companies discovers the live YC backend from the public page on each run.",
            "Primary filter: Spring 2025, agent / agentic pattern, team size ≤ 10.",
            "Structured cluster counts derived from 27 filtered companies (companies may span lanes).",
            "Use-case labels are inferred from company text, not official YC categories.",
        ],
        "This deck is chained from a structured analysis pack — not freehand slide writing.",
    )
    # 27 closing
    s_closing(
        prs, page(),
        [
            ("Lead with workflow-heavy use cases", "8 companies cluster around operational systems of action."),
            ("Own context, governance, or workflow position", "Durable clusters pair agents with data, infra, or embedded workflow."),
            ("Treat healthcare operations as a real wedge", "Smaller count, but strong pain and willingness to pay."),
            ("Avoid generic agent-platform messaging", "The cohort rewards explicit jobs and operators."),
        ],
    )

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUTPUT_PATH))
    print(f"Saved {OUTPUT_PATH} with {len(prs.slides._sldIdLst)} slides")


if __name__ == "__main__":
    main()
