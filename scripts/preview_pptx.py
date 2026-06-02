#!/usr/bin/env python3
"""Render a .pptx to PNG previews for visual QA (approximate, no auto-shrink).

Reads real shape geometry, fills, and text from the file and draws them with
matplotlib. Useful to catch text-overflow / overlap without PowerPoint.
"""
import argparse
import json
import sys
import textwrap
from pathlib import Path

import matplotlib

matplotlib.use("Agg")
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, Rectangle

from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu

EMU = 914400.0


def color_hex(shape):
    try:
        fill = shape.fill
        if fill.type is not None and fill.fore_color and fill.fore_color.type is not None:
            return "#" + str(fill.fore_color.rgb)
    except Exception:
        pass
    return None


def run_color(run):
    try:
        if run.font.color and run.font.color.type is not None:
            return "#" + str(run.font.color.rgb)
    except Exception:
        pass
    return "#1B2B3C"


def render(pptx_path: Path, out_dir: Path):
    prs = Presentation(str(pptx_path))
    W = prs.slide_width / EMU
    H = prs.slide_height / EMU
    out_dir.mkdir(parents=True, exist_ok=True)
    paths = []
    overflow_shapes = 0
    overflow_slides = []
    for idx, slide in enumerate(prs.slides, 1):
        fig, ax = plt.subplots(figsize=(W, H), dpi=110)
        ax.set_xlim(0, W)
        ax.set_ylim(0, H)
        ax.invert_yaxis()
        ax.axis("off")
        ax.add_patch(Rectangle((0, 0), W, H, facecolor="white", edgecolor="none"))
        for shape in slide.shapes:
            try:
                left = shape.left / EMU
                top = shape.top / EMU
                w = shape.width / EMU
                h = shape.height / EMU
            except Exception:
                continue
            if shape.shape_type == 13:  # picture
                ax.add_patch(Rectangle((left, top), w, h, facecolor="#dfe6ec", edgecolor="#b8c2cc"))
                ax.text(left + w / 2, top + h / 2, "[image]", ha="center", va="center", fontsize=8, color="#7C8B9A")
                continue
            fill = color_hex(shape)
            if fill and not shape.has_text_frame:
                ax.add_patch(FancyBboxPatch((left, top), w, h, boxstyle="round,pad=0,rounding_size=0.02",
                                            facecolor=fill, edgecolor="none"))
            elif fill and shape.has_text_frame and shape.text.strip() == "":
                ax.add_patch(FancyBboxPatch((left, top), w, h, boxstyle="round,pad=0,rounding_size=0.02",
                                            facecolor=fill, edgecolor="none"))
            if shape.has_text_frame and shape.text.strip():
                if fill:  # filled shape with text (badge/button)
                    ax.add_patch(FancyBboxPatch((left, top), w, h, boxstyle="round,pad=0,rounding_size=0.02",
                                                facecolor=fill, edgecolor="none"))
                tf = shape.text_frame
                # vertical anchor
                anchor = tf.vertical_anchor
                y = top + 0.04
                # collect paragraph lines
                rendered = []
                for para in tf.paragraphs:
                    runs = para.runs
                    if not runs:
                        continue
                    ptext = "".join(r.text for r in runs)
                    size = runs[0].font.size.pt if runs[0].font.size else 14
                    bold = bool(runs[0].font.bold)
                    col = run_color(runs[0])
                    align = para.alignment
                    # wrap by width: approx chars per inch at this pt size
                    char_w = size * 0.0085  # inches per char (rough)
                    max_chars = max(4, int(w / char_w))
                    wrapped = textwrap.wrap(ptext, max_chars) or [ptext]
                    sb = para.space_before.pt / 72.0 if para.space_before else 0
                    rendered.append((wrapped, size, bold, col, align, sb))
                total_h = sum((sb + len(wr) * (sz / 72.0) * 1.18) for wr, sz, _, _, _, sb in rendered)
                if anchor == 3:  # middle
                    y = top + max(0, (h - total_h) / 2)
                elif anchor == 4:  # bottom
                    y = top + max(0, h - total_h)
                overflow = total_h > h + 0.03
                for wrapped, size, bold, col, align, sb in rendered:
                    y += sb
                    for line in wrapped:
                        if align == PP_ALIGN.CENTER:
                            xx, ha = left + w / 2, "center"
                        elif align == PP_ALIGN.RIGHT:
                            xx, ha = left + w, "right"
                        else:
                            xx, ha = left, "left"
                        lh = (size / 72.0) * 1.18
                        ax.text(xx, y + lh * 0.5, line, ha=ha, va="center",
                                fontsize=size * 0.92, color=col, fontweight="bold" if bold else "normal")
                        y += lh
                if overflow:
                    overflow_shapes += 1
                    overflow_slides.append(idx)
                    ax.add_patch(Rectangle((left, top), w, h, facecolor="none",
                                           edgecolor="red", linewidth=1.5, linestyle="--"))
        out = out_dir / f"slide_{idx:02d}.png"
        fig.savefig(out, dpi=110, bbox_inches="tight", pad_inches=0)
        plt.close(fig)
        paths.append(out)
    # contact sheets of 9
    for start in range(0, len(paths), 9):
        chunk = paths[start:start + 9]
        fig, axes = plt.subplots(3, 3, figsize=(18, 10))
        for axx, pth in zip(axes.flat, chunk):
            img = plt.imread(pth)
            axx.imshow(img)
            axx.set_title(pth.stem, fontsize=9)
            axx.axis("off")
        for axx in axes.flat[len(chunk):]:
            axx.axis("off")
        sheet = out_dir / f"contact_{start//9 + 1}.png"
        fig.savefig(sheet, dpi=90, bbox_inches="tight")
        plt.close(fig)
        print("wrote", sheet)
    summary = {
        "pptx": str(pptx_path),
        "out_dir": str(out_dir),
        "slide_count": len(paths),
        "overflow_shape_count": overflow_shapes,
        "overflow_slides": sorted(set(overflow_slides)),
    }
    (out_dir / "qa-summary.json").write_text(json.dumps(summary, indent=2) + "\n", encoding="utf-8")
    print(f"Rendered {len(paths)} slides to {out_dir}")
    print(json.dumps(summary))
    return summary


if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Render PPTX previews and emit a QA summary.")
    parser.add_argument("pptx", nargs="?", default="docs/reports/yc-agent-companies-spring-2025-v2.pptx")
    parser.add_argument("--out-dir", default="docs/reports/_preview")
    args = parser.parse_args()
    render(Path(args.pptx), Path(args.out_dir))
