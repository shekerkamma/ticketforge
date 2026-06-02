#!/usr/bin/env python3
"""Build a branded YC agent companies deck from a structured analysis pack."""

from __future__ import annotations

import io
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


NAVY = RGBColor(0x0A, 0x16, 0x28)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
TEAL_DARK = RGBColor(0x00, 0x9B, 0x82)
CREAM = RGBColor(0xEE, 0xEC, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
GRID = RGBColor(0xD9, 0xDF, 0xE5)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)

TEMPLATE_PATH = Path("/mnt/c/Users/sheke/OneDrive/Desktop/Prasad_Agentic_AI_Use_Cases_Across_Industries.pptx")
DECK_PLAN_PATH = Path("analytics-comms/yc-agent-companies-spring-2025/deck-plan.json")
ANALYSIS_PATH = Path("analytics-comms/yc-agent-companies-spring-2025/analysis.json")
OUTPUT_PATH = Path("docs/reports/yc-agent-companies-spring-2025-branded.pptx")
TITLE_IMAGE_SLIDE_INDEX = 0
TITLE_IMAGE_SHAPE_INDEX = 2
BLANK_LAYOUT_INDEX = 6


def remove_all_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[idx]


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill if line is None else line
    if radius:
        shape.adjustments[0] = 0.15
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
    vertical=MSO_ANCHOR.TOP,
    fit=True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    frame.margin_left = Emu(35000)
    frame.margin_right = Emu(35000)
    frame.margin_top = Emu(20000)
    frame.margin_bottom = Emu(20000)
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    if fit:
        try:
            frame.fit_text(font_family=font_name, max_size=size, bold=bold)
        except Exception:
            pass
    return box


def add_body_lines(slide, lines: list[str], left, top, width, *, line_height: int, size: int, color) -> None:
    for idx, line in enumerate(lines):
        add_text(slide, line, left, top + Emu(line_height * idx), width, Emu(line_height), size=size, color=color)


def add_footer(slide, page_num: int) -> None:
    add_text(
        slide,
        "TicketForge | YC Agent Companies | June 2026",
        Emu(300000),
        Emu(6400000),
        Emu(4200000),
        Emu(180000),
        size=12,
        color=MUTED,
    )
    add_text(
        slide,
        str(page_num),
        Emu(11600000),
        Emu(6380000),
        Emu(220000),
        Emu(200000),
        size=16,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def add_title_slide(prs: Presentation, image_blob: bytes, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    slide.shapes.add_picture(io.BytesIO(image_blob), Emu(8076895), 0, Emu(4114800), SLIDE_H)
    words = spec["title"].split()
    if len(words) >= 3:
        line_one = " ".join(words[:2])
        line_two = " ".join(words[2:])
    elif len(words) == 2:
        line_one, line_two = words
    else:
        line_one, line_two = spec["title"], ""
    add_text(slide, line_one, Emu(731520), Emu(914400), Emu(6400800), Emu(700000), size=28, color=WHITE)
    add_text(slide, line_two, Emu(731520), Emu(1550000), Emu(6400800), Emu(950000), size=42, color=WHITE, bold=True)
    add_rect(slide, Emu(731520), Emu(2700000), Emu(1371600), Emu(36576), TEAL)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(3000000), Emu(6400800), Emu(400000), size=22, color=CREAM, bold=True)
    add_text(slide, spec["strapline"], Emu(731520), Emu(3450000), Emu(6500000), Emu(900000), size=19, color=CREAM)
    chips = ["Workflow", "Customer Ops", "Healthcare", "Tooling", "Infra"]
    x = Emu(731520)
    for chip in chips:
        add_rect(slide, x, Emu(4850000), Emu(920000), Emu(320000), TEAL_DARK, radius=True)
        add_text(slide, chip, x, Emu(4895000), Emu(920000), Emu(220000), size=12, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        x += Emu(1030000)


def add_agenda_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(7436815), Emu(457200), size=26, color=NAVY, bold=True)
    top = Emu(960120)
    for section in spec["sections"]:
        add_rect(slide, Emu(365760), top, Emu(457200), Emu(411480), NAVY)
        add_text(slide, section["number"], Emu(365760), top + Emu(70000), Emu(457200), Emu(274320), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, section["title"], Emu(960120), top + Emu(25000), Emu(4572000), Emu(220000), size=18, color=NAVY, bold=True)
        add_text(slide, section["detail"], Emu(960120), top + Emu(245000), Emu(5943600), Emu(220000), size=12, color=MUTED)
        top += Emu(570000)
    add_footer(slide, page_num)


def add_section_divider_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_text(slide, spec["section"], Emu(731520), Emu(1400000), Emu(2500000), Emu(220000), size=11, color=TEAL, bold=True)
    add_text(slide, spec["title"], Emu(731520), Emu(1850000), Emu(7000000), Emu(700000), size=28, color=WHITE, bold=True)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(2650000), Emu(7600000), Emu(700000), size=17, color=CREAM)
    add_rect(slide, Emu(731520), Emu(3550000), Emu(1371600), Emu(36576), TEAL)
    add_footer(slide, page_num)


def add_summary_cards_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_text(slide, spec["title"], Emu(731520), Emu(650000), Emu(9000000), Emu(600000), size=30, color=NAVY, bold=True)
    add_rect(slide, Emu(731520), Emu(1250000), Emu(1371600), Emu(36576), TEAL)

    top = Emu(1800000)
    for card in spec["cards"]:
        add_rect(slide, Emu(731520), top, Emu(10600000), Emu(1050000), WHITE, line=GRID, radius=True)
        add_rect(slide, Emu(900000), top + Emu(150000), Emu(520000), Emu(520000), NAVY, radius=True)
        add_text(slide, card["number"], Emu(900000), top + Emu(190000), Emu(520000), Emu(280000), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, card["title"], Emu(1600000), top + Emu(170000), Emu(6200000), Emu(260000), size=20, color=NAVY, bold=True)
        add_text(slide, card["body"], Emu(1600000), top + Emu(450000), Emu(9300000), Emu(420000), size=16, color=INK)
        top += Emu(1200000)
    add_footer(slide, page_num)


def add_kpi_grid_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    add_text(slide, spec["title"], Emu(731520), Emu(500000), Emu(9000000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(980000), Emu(10000000), Emu(320000), size=16, color=NAVY)

    lefts = [Emu(731520), Emu(3350000), Emu(5970000), Emu(8590000)]
    for kpi, left in zip(spec["kpis"], lefts):
        add_rect(slide, left, Emu(1800000), Emu(2200000), Emu(1600000), NAVY, radius=True)
        add_text(slide, kpi["value"], left, Emu(2050000), Emu(2200000), Emu(460000), size=28, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, kpi["label"], left + Emu(150000), Emu(2550000), Emu(1900000), Emu(260000), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, kpi["note"], left + Emu(150000), Emu(2920000), Emu(1900000), Emu(420000), size=13, color=CREAM, align=PP_ALIGN.CENTER)

    add_rect(slide, Emu(731520), Emu(4150000), Emu(10600000), Emu(1200000), CREAM, radius=True)
    add_text(
        slide,
        spec["takeaway"],
        Emu(1050000),
        Emu(4520000),
        Emu(10000000),
        Emu(450000),
        size=20,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
        vertical=MSO_ANCHOR.MIDDLE,
    )
    add_footer(slide, page_num)


def add_bar_chart_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_text(slide, spec["title"], Emu(731520), Emu(500000), Emu(9500000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(980000), Emu(9000000), Emu(260000), size=16, color=MUTED)

    max_value = max(item["value"] for item in spec["series"]) or 1
    top = Emu(1650000)
    for item in spec["series"]:
        label = item["label"]
        value = item["value"]
        width = int(5200000 * (value / max_value))
        fill = TEAL if value >= 6 else TEAL_DARK
        add_text(slide, label, Emu(900000), top, Emu(2900000), Emu(250000), size=17, color=INK, bold=True)
        add_rect(slide, Emu(3900000), top + Emu(20000), Emu(5200000), Emu(210000), CREAM, line=CREAM, radius=True)
        add_rect(slide, Emu(3900000), top + Emu(20000), Emu(width), Emu(210000), fill, line=fill, radius=True)
        add_text(slide, str(value), Emu(9300000), top - Emu(10000), Emu(350000), Emu(250000), size=17, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
        top += Emu(520000)

    add_text(slide, spec["takeaway"], Emu(900000), Emu(5650000), Emu(10000000), Emu(350000), size=18, color=NAVY, bold=True)
    add_footer(slide, page_num)


def add_team_distribution_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(8000000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(365760), Emu(700000), Emu(7600000), Emu(240000), size=12, color=MUTED)
    max_value = max(item["count"] for item in spec["series"]) or 1
    lefts = [Emu(900000), Emu(3000000), Emu(5100000), Emu(7200000), Emu(9300000)]
    for item, left in zip(spec["series"], lefts):
        height = int(1900000 * (item["count"] / max_value))
        add_rect(slide, left, Emu(3900000) - Emu(height), Emu(1200000), Emu(height), TEAL if item["count"] >= max_value else NAVY, radius=True)
        add_text(slide, str(item["count"]), left, Emu(3950000) - Emu(height), Emu(1200000), Emu(260000), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, item["label"], left, Emu(4100000), Emu(1200000), Emu(220000), size=15, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, spec["takeaway"], Emu(900000), Emu(5000000), Emu(9800000), Emu(500000), size=18, color=INK, bold=True, align=PP_ALIGN.CENTER)
    add_footer(slide, page_num)


def add_comparison_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(7436815), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(slide, "Previous expectation vs. what the YC cohort actually reveals", Emu(365760), Emu(680000), Emu(7000000), Emu(220000), size=12, color=MUTED)

    add_rect(slide, Emu(365760), Emu(1100000), Emu(3300000), Emu(420000), RGBColor(0xF5, 0x47, 0x76))
    add_rect(slide, Emu(4300000), Emu(1100000), Emu(3300000), Emu(420000), RGBColor(0x2D, 0xC4, 0x8D))
    add_text(slide, spec["left_title"], Emu(365760), Emu(1180000), Emu(3300000), Emu(220000), size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, spec["right_title"], Emu(4300000), Emu(1180000), Emu(3300000), Emu(220000), size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

    top = Emu(1620000)
    for left_item, right_item in zip(spec["left_points"], spec["right_points"]):
        add_rect(slide, Emu(365760), top, Emu(3300000), Emu(470000), RGBColor(0xFB, 0xF0, 0xF2), line=RGBColor(0xF5, 0xD7, 0xDF), radius=True)
        add_rect(slide, Emu(4300000), top, Emu(3300000), Emu(470000), RGBColor(0xEE, 0xFA, 0xF5), line=RGBColor(0xCC, 0xEE, 0xDE), radius=True)
        add_text(slide, f"×  {left_item}", Emu(520000), top + Emu(90000), Emu(2900000), Emu(260000), size=14, color=INK)
        add_text(slide, f"✓  {right_item}", Emu(4450000), top + Emu(90000), Emu(2900000), Emu(260000), size=14, color=INK)
        top += Emu(560000)
    add_footer(slide, page_num)


def add_structured_table_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(7436815), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(365760), Emu(680000), Emu(7600000), Emu(220000), size=12, color=MUTED)

    clusters = spec["clusters"][:4]
    headers = [cluster["name"].title() for cluster in clusters]
    col_lefts = [Emu(2400000), Emu(4550000), Emu(6700000), Emu(8850000)]
    add_text(slide, "", Emu(365760), Emu(1300000), Emu(1800000), Emu(200000), size=12, color=NAVY, bold=True)
    for header, left in zip(headers, col_lefts):
        add_text(slide, header, left, Emu(1240000), Emu(1800000), Emu(240000), size=13, color=NAVY, bold=True, align=PP_ALIGN.CENTER)

    row_labels = ["Buyer", "Business job", "Value prop", "Why now"]
    row_tops = [Emu(1650000), Emu(2250000), Emu(3150000), Emu(4050000)]
    fields = ["buyer", "business_job", "value_prop", "why_now"]
    for label, top in zip(row_labels, row_tops):
        add_rect(slide, Emu(365760), top, Emu(1500000), Emu(420000), TEAL if label != "Business job" else NAVY)
        add_text(slide, label, Emu(420000), top + Emu(80000), Emu(1350000), Emu(220000), size=13, color=WHITE, bold=True)
    for cluster, left in zip(clusters, col_lefts):
        for field, top in zip(fields, row_tops):
            add_rect(slide, left, top, Emu(1900000), Emu(420000), RGBColor(0xF3, 0xF6, 0xF8), line=GRID)
            add_text(slide, cluster[field], left + Emu(120000), top + Emu(70000), Emu(1660000), Emu(300000), size=11, color=INK)
    add_footer(slide, page_num)


def add_cluster_spotlight_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_text(slide, spec["title"], Emu(731520), Emu(420000), Emu(8000000), Emu(500000), size=28, color=NAVY, bold=True)
    add_rect(slide, Emu(731520), Emu(930000), Emu(1371600), Emu(36576), TEAL)
    add_text(slide, f"Buyer: {spec['buyer']}", Emu(731520), Emu(1180000), Emu(4600000), Emu(240000), size=13, color=MUTED)
    add_text(slide, f"Workflow owner: {spec['workflow_owner']}", Emu(731520), Emu(1450000), Emu(4600000), Emu(260000), size=13, color=MUTED)
    add_rect(slide, Emu(731520), Emu(1850000), Emu(4900000), Emu(1500000), WHITE, line=GRID, radius=True)
    add_text(slide, "Business job", Emu(900000), Emu(2010000), Emu(1600000), Emu(220000), size=15, color=NAVY, bold=True)
    add_text(slide, spec["business_job"], Emu(900000), Emu(2310000), Emu(4500000), Emu(360000), size=14, color=INK)
    add_text(slide, "Value proposition", Emu(900000), Emu(2740000), Emu(1800000), Emu(220000), size=15, color=NAVY, bold=True)
    add_text(slide, spec["value_prop"], Emu(900000), Emu(3040000), Emu(4500000), Emu(300000), size=13, color=INK)
    add_text(slide, "Why now", Emu(731520), Emu(3550000), Emu(1200000), Emu(220000), size=15, color=NAVY, bold=True)
    add_text(slide, spec["why_now"], Emu(900000), Emu(3850000), Emu(4800000), Emu(300000), size=13, color=INK)
    add_text(slide, "Recommendation", Emu(731520), Emu(4300000), Emu(1500000), Emu(220000), size=15, color=NAVY, bold=True)
    add_text(slide, spec["recommendation"], Emu(900000), Emu(4600000), Emu(4800000), Emu(300000), size=13, color=INK)

    add_rect(slide, Emu(6100000), Emu(1450000), Emu(4700000), Emu(4000000), NAVY, radius=True)
    add_text(slide, "Source-backed examples", Emu(6350000), Emu(1650000), Emu(2500000), Emu(220000), size=16, color=WHITE, bold=True)
    top = Emu(2050000)
    for example in spec["examples"][:4]:
        add_rect(slide, Emu(6350000), top, Emu(4200000), Emu(700000), WHITE, line=WHITE, radius=True)
        add_text(slide, example["name"], Emu(6550000), top + Emu(70000), Emu(1500000), Emu(220000), size=14, color=NAVY, bold=True)
        add_text(slide, example["proof_line"], Emu(6550000), top + Emu(280000), Emu(3800000), Emu(320000), size=11, color=INK)
        top += Emu(850000)
    add_footer(slide, page_num)


def add_vertical_table_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_text(slide, spec["title"], Emu(731520), Emu(500000), Emu(9800000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(980000), Emu(9800000), Emu(260000), size=16, color=MUTED)

    add_rect(slide, Emu(731520), Emu(1600000), Emu(10600000), Emu(2300000), CREAM, line=GRID, radius=True)
    headers = ["Lane", "Count", "Business job", "Evidence"]
    xs = [Emu(900000), Emu(2800000), Emu(3600000), Emu(8400000)]
    widths = [Emu(1700000), Emu(600000), Emu(4500000), Emu(2400000)]
    for header, left, width in zip(headers, xs, widths):
        add_text(slide, header, left, Emu(1760000), width, Emu(220000), size=13, color=NAVY, bold=True)

    top = Emu(2100000)
    for cluster in spec["clusters"][:4]:
        add_text(slide, cluster["name"].title(), xs[0], top, widths[0], Emu(240000), size=13, color=INK, bold=True)
        add_text(slide, str(cluster["count"]), xs[1], top, widths[1], Emu(240000), size=13, color=NAVY, bold=True)
        add_text(slide, cluster["business_job"], xs[2], top, widths[2], Emu(420000), size=12, color=INK)
        add_text(slide, ", ".join(cluster["examples"][:3]), xs[3], top, widths[3], Emu(240000), size=12, color=MUTED)
        top += Emu(460000)

    add_rect(slide, Emu(731520), Emu(4300000), Emu(10600000), Emu(1300000), NAVY, radius=True)
    add_text(slide, "Representative companies across lanes", Emu(950000), Emu(4470000), Emu(3000000), Emu(220000), size=16, color=WHITE, bold=True)
    row_top = Emu(4800000)
    for company in spec["companies"][:6]:
        add_text(
            slide,
            f"{company['name']} | {company['cluster']} | {company['what_they_do']}",
            Emu(950000),
            row_top,
            Emu(9500000),
            Emu(220000),
            size=12,
            color=CREAM,
        )
        row_top += Emu(210000)
    add_footer(slide, page_num)


def add_company_grid_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(7800000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(365760), Emu(700000), Emu(7600000), Emu(220000), size=12, color=MUTED)
    positions = [
        (Emu(365760), Emu(1300000)),
        (Emu(6200000), Emu(1300000)),
        (Emu(365760), Emu(3600000)),
        (Emu(6200000), Emu(3600000)),
    ]
    for company, (left, top) in zip(spec["companies"], positions):
        add_rect(slide, left, top, Emu(5000000), Emu(1800000), CREAM, line=GRID, radius=True)
        add_rect(slide, left + Emu(170000), top + Emu(150000), Emu(1100000), Emu(220000), TEAL_DARK, radius=True)
        add_text(slide, company["cluster"].title(), left + Emu(170000), top + Emu(165000), Emu(1100000), Emu(150000), size=11, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, company["name"], left + Emu(170000), top + Emu(470000), Emu(3000000), Emu(240000), size=18, color=NAVY, bold=True)
        add_text(slide, company["proof_line"], left + Emu(170000), top + Emu(820000), Emu(4500000), Emu(700000), size=12, color=INK)
    add_footer(slide, page_num)


def add_case_study_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    add_text(slide, spec["title"], Emu(365760), Emu(182880), Emu(9000000), Emu(502920), size=28, color=NAVY, bold=True)
    add_rect(slide, Emu(365760), Emu(658368), Emu(1828800), Emu(27432), NAVY)
    add_text(slide, spec["match_title"].upper(), Emu(365760), Emu(777240), Emu(2200000), Emu(200000), size=16, color=NAVY, bold=True)

    left_lines = [f"• {item['name']} — {item.get('proof_line') or item.get('one_liner')}" for item in spec["matches"][:4]]
    add_body_lines(slide, left_lines, Emu(365760), Emu(960120), Emu(3520440), line_height=300000, size=14, color=INK)

    add_text(slide, spec["implication_title"].upper(), Emu(4023360), Emu(777240), Emu(2200000), Emu(200000), size=16, color=NAVY, bold=True)
    right_lines = [f"• {line}" for line in spec["implications"]]
    add_body_lines(slide, right_lines, Emu(4023360), Emu(960120), Emu(3520440), line_height=300000, size=14, color=INK)

    lefts = [Emu(365760), Emu(2240279), Emu(4114799)]
    for metric, left in zip(spec["kpis"], lefts):
        add_rect(slide, left, Emu(2500000), Emu(1645920), Emu(822960), NAVY, radius=True)
        add_text(slide, metric["value"], left, Emu(2550000), Emu(1645920), Emu(411480), size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, metric["label"], left + Emu(100000), Emu(2920000), Emu(1445920), Emu(320040), size=12, color=WHITE, align=PP_ALIGN.CENTER)

    add_text(
        slide,
        "The key question is whether the product owns enough workflow to become system-of-action software.",
        Emu(365760),
        Emu(3900000),
        Emu(8000000),
        Emu(500000),
        size=18,
        color=NAVY,
        bold=True,
    )
    add_footer(slide, page_num)


def add_implication_bullets_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(7600000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(365760), Emu(700000), Emu(8200000), Emu(260000), size=13, color=MUTED)
    top = Emu(1350000)
    for idx, bullet in enumerate(spec["bullets"], start=1):
        add_rect(slide, Emu(500000), top, Emu(520000), Emu(520000), NAVY, radius=True)
        add_text(slide, f"{idx:02d}", Emu(500000), top + Emu(120000), Emu(520000), Emu(180000), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, Emu(1200000), top, Emu(9300000), Emu(620000), RGBColor(0xF3, 0xF6, 0xF8), line=GRID, radius=True)
        add_text(slide, bullet, Emu(1450000), top + Emu(140000), Emu(8800000), Emu(340000), size=15, color=INK)
        top += Emu(760000)
    add_footer(slide, page_num)


def add_recommendation_grid_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_text(slide, spec["title"], Emu(731520), Emu(500000), Emu(9000000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(980000), Emu(9000000), Emu(260000), size=16, color=MUTED)

    positions = [
        (Emu(731520), Emu(1700000)),
        (Emu(6100000), Emu(1700000)),
        (Emu(731520), Emu(3820000)),
        (Emu(6100000), Emu(3820000)),
    ]
    for idx, (rec, (left, top)) in enumerate(zip(spec["recommendations"], positions), start=1):
        add_rect(slide, left, top, Emu(4700000), Emu(1600000), WHITE, line=GRID, radius=True)
        add_rect(slide, left + Emu(180000), top + Emu(160000), Emu(520000), Emu(520000), TEAL, radius=True)
        add_text(slide, f"{idx:02d}", left + Emu(180000), top + Emu(220000), Emu(520000), Emu(220000), size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, rec["title"], left + Emu(900000), top + Emu(180000), Emu(3400000), Emu(260000), size=18, color=NAVY, bold=True)
        add_text(slide, rec["why"], left + Emu(900000), top + Emu(560000), Emu(3400000), Emu(520000), size=14, color=INK)
    add_footer(slide, page_num)


def add_opportunity_ladder_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_text(slide, spec["title"], Emu(731520), Emu(500000), Emu(8000000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(731520), Emu(980000), Emu(9000000), Emu(240000), size=13, color=MUTED)
    top = Emu(1650000)
    widths = [Emu(1500000), Emu(3000000), Emu(4700000)]
    for idx, item in enumerate(spec["items"]):
        add_rect(slide, Emu(900000), top, Emu(1200000), Emu(520000), TEAL if idx < 3 else NAVY, radius=True)
        add_text(slide, item["label"], Emu(900000), top + Emu(140000), Emu(1200000), Emu(180000), size=14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_rect(slide, Emu(2400000), top, Emu(3400000), Emu(520000), WHITE, line=GRID, radius=True)
        add_text(slide, item["value"], Emu(2580000), top + Emu(120000), Emu(3000000), Emu(220000), size=16, color=NAVY, bold=True)
        add_rect(slide, Emu(6100000), top, Emu(4300000), Emu(520000), WHITE, line=GRID, radius=True)
        add_text(slide, item["detail"], Emu(6280000), top + Emu(120000), Emu(3900000), Emu(220000), size=12, color=INK)
        top += Emu(760000)
    add_footer(slide, page_num)


def add_methodology_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_text(slide, spec["title"], Emu(731520), Emu(500000), Emu(9000000), Emu(500000), size=30, color=NAVY, bold=True)
    add_rect(slide, Emu(731520), Emu(1250000), Emu(1371600), Emu(36576), TEAL)
    lines = [f"• {bullet}" for bullet in spec["bullets"]]
    add_body_lines(slide, lines, Emu(900000), Emu(1700000), Emu(10000000), line_height=330000, size=16, color=INK)
    add_rect(slide, Emu(900000), Emu(4150000), Emu(10300000), Emu(1100000), CREAM, line=GRID, radius=True)
    add_text(
        slide,
        "This deck is chained from a structured analysis pack, not from freehand slide writing.",
        Emu(1100000),
        Emu(4550000),
        Emu(9700000),
        Emu(300000),
        size=18,
        color=NAVY,
        bold=True,
        align=PP_ALIGN.CENTER,
    )
    add_footer(slide, page_num)


def add_appendix_sources_slide(prs: Presentation, page_num: int, spec: dict[str, Any]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[BLANK_LAYOUT_INDEX])
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(slide, spec["title"], Emu(365760), Emu(274320), Emu(7600000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(slide, spec["subtitle"], Emu(365760), Emu(700000), Emu(7600000), Emu(220000), size=12, color=MUTED)
    top = Emu(1300000)
    for company in spec["companies"][:8]:
        add_rect(slide, Emu(365760), top, Emu(10600000), Emu(460000), RGBColor(0xF7, 0xF9, 0xFB), line=GRID)
        add_text(slide, company["name"], Emu(520000), top + Emu(100000), Emu(1900000), Emu(180000), size=13, color=NAVY, bold=True)
        add_text(slide, company["proof_line"], Emu(2400000), top + Emu(90000), Emu(8200000), Emu(220000), size=11, color=INK)
        top += Emu(520000)
    add_footer(slide, page_num)


SLIDE_BUILDERS = {
    "agenda": add_agenda_slide,
    "section_divider": add_section_divider_slide,
    "summary_cards": add_summary_cards_slide,
    "kpi_grid": add_kpi_grid_slide,
    "team_distribution": add_team_distribution_slide,
    "bar_chart": add_bar_chart_slide,
    "comparison": add_comparison_slide,
    "structured_table": add_structured_table_slide,
    "cluster_spotlight": add_cluster_spotlight_slide,
    "vertical_table": add_vertical_table_slide,
    "company_grid": add_company_grid_slide,
    "case_study": add_case_study_slide,
    "implication_bullets": add_implication_bullets_slide,
    "recommendation_grid": add_recommendation_grid_slide,
    "opportunity_ladder": add_opportunity_ladder_slide,
    "methodology": add_methodology_slide,
    "appendix_sources": add_appendix_sources_slide,
}


def render_spec_from_slide(slide: dict[str, Any]) -> dict[str, Any]:
    if "render_spec" in slide:
        return slide["render_spec"]
    spec = {
        "type": slide.get("layout", slide.get("slide_type", "custom")).replace("-", "_"),
        "title": slide.get("title", ""),
    }
    if slide.get("content_blocks"):
        spec["content_blocks"] = slide["content_blocks"]
    return spec


def load_slide_specs() -> list[dict[str, Any]]:
    if DECK_PLAN_PATH.exists():
        deck_plan = json.loads(DECK_PLAN_PATH.read_text(encoding="utf-8"))
        return [render_spec_from_slide(slide) for slide in deck_plan["slides"]]
    if ANALYSIS_PATH.exists():
        analysis = json.loads(ANALYSIS_PATH.read_text(encoding="utf-8"))
        return analysis["deck_plan"]
    raise SystemExit(f"Missing deck inputs: {DECK_PLAN_PATH} or {ANALYSIS_PATH}")


def main() -> None:
    slide_specs = load_slide_specs()

    prs = Presentation(str(TEMPLATE_PATH))
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    title_image = prs.slides[TITLE_IMAGE_SLIDE_INDEX].shapes[TITLE_IMAGE_SHAPE_INDEX].image.blob
    remove_all_slides(prs)

    add_title_slide(prs, title_image, slide_specs[0])
    for page_num, spec in enumerate(slide_specs[1:], start=2):
        builder = SLIDE_BUILDERS[spec["type"]]
        builder(prs, page_num, spec)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    temp_path = OUTPUT_PATH.with_suffix(".tmp.pptx")
    prs.save(str(temp_path))
    temp_path.replace(OUTPUT_PATH)


if __name__ == "__main__":
    main()
