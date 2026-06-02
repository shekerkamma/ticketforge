#!/usr/bin/env python3
"""Render any deck-plan.json into a branded PPTX."""

from __future__ import annotations

import argparse
import io
import json
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Pt


NAVY = RGBColor(0x0A, 0x16, 0x28)
TEAL = RGBColor(0x00, 0xC9, 0xA7)
TEAL_DARK = RGBColor(0x00, 0x9B, 0x82)
CREAM = RGBColor(0xEE, 0xEC, 0xE1)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
INK = RGBColor(0x1B, 0x2B, 0x3C)
MUTED = RGBColor(0x5B, 0x6B, 0x7C)
GRID = RGBColor(0xD9, 0xDF, 0xE5)
ACCENT_RED = RGBColor(0xF5, 0x47, 0x76)
ACCENT_GREEN = RGBColor(0x2D, 0xC4, 0x8D)

SLIDE_W = Emu(12192000)
SLIDE_H = Emu(6858000)
DEFAULT_TEMPLATE = Path("/mnt/c/Users/sheke/OneDrive/Desktop/Prasad_Agentic_AI_Use_Cases_Across_Industries.pptx")
TITLE_IMAGE_SLIDE_INDEX = 0
TITLE_IMAGE_SHAPE_INDEX = 2
BLANK_LAYOUT_INDEX = 6


def remove_all_slides(prs: Presentation) -> None:
    for idx in range(len(prs.slides) - 1, -1, -1):
        r_id = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(r_id)
        del prs.slides._sldIdLst[idx]


def add_rect(slide, left, top, width, height, fill, line=None, radius=False):
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if radius else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill if line is None else line
    if radius:
        shape.adjustments[0] = 0.15
    return shape


def add_text(
    slide,
    text,
    left,
    top,
    width,
    height,
    *,
    size,
    color,
    bold=False,
    align=PP_ALIGN.LEFT,
    font_name="Calibri",
    vertical=MSO_ANCHOR.TOP,
    fit=True,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.vertical_anchor = vertical
    frame.margin_left = Emu(35000)
    frame.margin_right = Emu(35000)
    frame.margin_top = Emu(20000)
    frame.margin_bottom = Emu(20000)
    para = frame.paragraphs[0]
    para.alignment = align
    run = para.add_run()
    run.text = text
    font = run.font
    font.name = font_name
    font.size = Pt(size)
    font.bold = bold
    font.color.rgb = color
    if fit:
        try:
            frame.fit_text(font_family=font_name, max_size=size, bold=bold)
        except Exception:
            pass
    return box


def add_footer(slide, page_num: int, footer_text: str) -> None:
    add_text(
        slide,
        footer_text,
        Emu(300000),
        Emu(6400000),
        Emu(4200000),
        Emu(180000),
        size=12,
        color=MUTED,
    )
    add_text(
        slide,
        str(page_num),
        Emu(11600000),
        Emu(6380000),
        Emu(220000),
        Emu(200000),
        size=16,
        color=MUTED,
        align=PP_ALIGN.RIGHT,
    )


def block_body(slide: dict[str, Any], kind: str | None = None, label: str | None = None) -> list[dict[str, str]]:
    blocks = slide.get("content_blocks", [])
    out = []
    for block in blocks:
        if kind is not None and block.get("kind") != kind:
            continue
        if label is not None and block.get("label") != label:
            continue
        out.append(block)
    return out


def first_body(slide: dict[str, Any], kind: str | None = None, label: str | None = None, default: str = "") -> str:
    blocks = block_body(slide, kind=kind, label=label)
    return blocks[0]["body"] if blocks else default


def first_label(slide: dict[str, Any], kind: str | None = None, default: str = "") -> str:
    blocks = block_body(slide, kind=kind)
    return blocks[0]["label"] if blocks else default


def get_title_image(template_path: Path) -> bytes | None:
    if not template_path.exists():
        return None
    prs = Presentation(str(template_path))
    try:
        return prs.slides[TITLE_IMAGE_SLIDE_INDEX].shapes[TITLE_IMAGE_SHAPE_INDEX].image.blob
    except Exception:
        return None


def build_presentation(template_path: Path) -> Presentation:
    prs = Presentation(str(template_path)) if template_path.exists() else Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    remove_all_slides(prs)
    return prs


def slide_layout(prs: Presentation):
    try:
        return prs.slide_layouts[BLANK_LAYOUT_INDEX]
    except Exception:
        return prs.slide_layouts[0]


def rank_value(text: str) -> float:
    text = text.strip().lower()
    named = {
        "highest": 5.0,
        "high": 4.5,
        "medium-high": 3.8,
        "medium": 3.0,
        "medium-low": 2.2,
        "low": 1.5,
        "lowest": 1.0,
    }
    if text in named:
        return named[text]
    try:
        return float(text)
    except ValueError:
        return 0.0


def render_hero(prs: Presentation, image_blob: bytes | None, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    if image_blob:
        try:
            ppt_slide.shapes.add_picture(io.BytesIO(image_blob), Emu(8076895), 0, Emu(4114800), SLIDE_H)
        except Exception:
            pass
    title = slide["title"]
    subtitle = first_body(slide, label="Subtitle", default=slide.get("objective", ""))
    strapline = first_body(slide, label="Strapline", default="")
    words = title.split()
    line_one = " ".join(words[:2]) if len(words) >= 2 else title
    line_two = " ".join(words[2:]) if len(words) > 2 else ""
    add_text(ppt_slide, line_one, Emu(731520), Emu(914400), Emu(6400800), Emu(700000), size=28, color=WHITE)
    add_text(ppt_slide, line_two or title, Emu(731520), Emu(1550000), Emu(6400800), Emu(950000), size=42, color=WHITE, bold=True)
    add_rect(ppt_slide, Emu(731520), Emu(2700000), Emu(1371600), Emu(36576), TEAL)
    if subtitle:
        add_text(ppt_slide, subtitle, Emu(731520), Emu(3000000), Emu(6400800), Emu(460000), size=22, color=CREAM, bold=True)
    if strapline:
        add_text(ppt_slide, strapline, Emu(731520), Emu(3500000), Emu(6500000), Emu(900000), size=18, color=CREAM)
    add_footer(ppt_slide, page_num, footer_text)


def render_agenda(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(ppt_slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(ppt_slide, slide["title"], Emu(365760), Emu(274320), Emu(7436815), Emu(457200), size=26, color=NAVY, bold=True)
    top = Emu(1000000)
    for idx, block in enumerate(block_body(slide), start=1):
        add_rect(ppt_slide, Emu(365760), top, Emu(457200), Emu(411480), NAVY)
        add_text(ppt_slide, f"{idx:02d}", Emu(365760), top + Emu(70000), Emu(457200), Emu(274320), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(ppt_slide, block["label"], Emu(960120), top + Emu(25000), Emu(4572000), Emu(220000), size=18, color=NAVY, bold=True)
        add_text(ppt_slide, block["body"], Emu(960120), top + Emu(245000), Emu(6500000), Emu(220000), size=12, color=MUTED)
        top += Emu(570000)
    add_footer(ppt_slide, page_num, footer_text)


def render_section_divider(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, NAVY)
    add_text(ppt_slide, slide.get("section", "SECTION"), Emu(731520), Emu(1400000), Emu(2500000), Emu(220000), size=11, color=TEAL, bold=True)
    add_text(ppt_slide, slide["title"], Emu(731520), Emu(1850000), Emu(7000000), Emu(700000), size=28, color=WHITE, bold=True)
    add_text(ppt_slide, first_body(slide, kind="callout", default=slide.get("objective", "")), Emu(731520), Emu(2650000), Emu(7600000), Emu(700000), size=17, color=CREAM)
    add_rect(ppt_slide, Emu(731520), Emu(3550000), Emu(1371600), Emu(36576), TEAL)
    add_footer(ppt_slide, page_num, footer_text)


def render_card_grid(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str, dark=False) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, TEAL if dark else CREAM)
    title_color = NAVY if dark else NAVY
    add_text(ppt_slide, slide["title"], Emu(731520), Emu(650000), Emu(9000000), Emu(600000), size=30, color=title_color, bold=True)
    subtitle = first_body(slide, label="Subtitle", default="")
    if subtitle:
        add_text(ppt_slide, subtitle, Emu(731520), Emu(1060000), Emu(9300000), Emu(260000), size=15, color=MUTED)
    cards = block_body(slide)
    positions = [
        (Emu(731520), Emu(1700000)),
        (Emu(6100000), Emu(1700000)),
        (Emu(731520), Emu(3820000)),
        (Emu(6100000), Emu(3820000)),
    ]
    for idx, (block, (left, top)) in enumerate(zip(cards, positions), start=1):
        add_rect(ppt_slide, left, top, Emu(4700000), Emu(1600000), WHITE, line=GRID, radius=True)
        add_rect(ppt_slide, left + Emu(180000), top + Emu(160000), Emu(520000), Emu(520000), TEAL, radius=True)
        add_text(ppt_slide, f"{idx:02d}", left + Emu(180000), top + Emu(220000), Emu(520000), Emu(220000), size=17, color=NAVY, bold=True, align=PP_ALIGN.CENTER)
        add_text(ppt_slide, block["label"], left + Emu(900000), top + Emu(180000), Emu(3400000), Emu(260000), size=18, color=NAVY, bold=True)
        add_text(ppt_slide, block["body"], left + Emu(900000), top + Emu(560000), Emu(3400000), Emu(520000), size=14, color=INK)
    add_footer(ppt_slide, page_num, footer_text)


def render_metric_grid(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    add_text(ppt_slide, slide["title"], Emu(731520), Emu(500000), Emu(9000000), Emu(500000), size=30, color=NAVY, bold=True)
    metrics = block_body(slide, kind="metric")
    lefts = [Emu(731520), Emu(3350000), Emu(5970000), Emu(8590000)]
    for block, left in zip(metrics[:4], lefts):
        parts = [part.strip() for part in block["body"].split("|")]
        value = parts[0] if parts else block["body"]
        note = parts[1] if len(parts) > 1 else ""
        add_rect(ppt_slide, left, Emu(1800000), Emu(2200000), Emu(1600000), NAVY, radius=True)
        add_text(ppt_slide, value, left, Emu(2050000), Emu(2200000), Emu(460000), size=28, color=TEAL, bold=True, align=PP_ALIGN.CENTER)
        add_text(ppt_slide, block["label"], left + Emu(150000), Emu(2550000), Emu(1900000), Emu(260000), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_text(ppt_slide, note, left + Emu(150000), Emu(2920000), Emu(1900000), Emu(420000), size=13, color=CREAM, align=PP_ALIGN.CENTER)
    takeaway = first_body(slide, label="Takeaway", default=slide.get("objective", ""))
    if takeaway:
        add_rect(ppt_slide, Emu(731520), Emu(4150000), Emu(10600000), Emu(1200000), CREAM, radius=True)
        add_text(ppt_slide, takeaway, Emu(1050000), Emu(4520000), Emu(10000000), Emu(450000), size=20, color=NAVY, bold=True, align=PP_ALIGN.CENTER, vertical=MSO_ANCHOR.MIDDLE)
    add_footer(ppt_slide, page_num, footer_text)


def render_chart_like(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_text(ppt_slide, slide["title"], Emu(731520), Emu(500000), Emu(9500000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(ppt_slide, slide.get("objective", ""), Emu(731520), Emu(980000), Emu(9000000), Emu(260000), size=16, color=MUTED)
    items = block_body(slide, kind="chart-brief")
    max_value = max((rank_value(item["body"]) for item in items), default=1.0) or 1.0
    top = Emu(1650000)
    for item in items[:8]:
        value = rank_value(item["body"])
        width = int(5200000 * ((value or 1.0) / max_value))
        fill = TEAL if value >= max_value * 0.8 else TEAL_DARK
        add_text(ppt_slide, item["label"], Emu(900000), top, Emu(2900000), Emu(250000), size=17, color=INK, bold=True)
        add_rect(ppt_slide, Emu(3900000), top + Emu(20000), Emu(5200000), Emu(210000), CREAM, line=CREAM, radius=True)
        add_rect(ppt_slide, Emu(3900000), top + Emu(20000), Emu(max(width, 250000)), Emu(210000), fill, line=fill, radius=True)
        add_text(ppt_slide, item["body"], Emu(9300000), top - Emu(10000), Emu(900000), Emu(250000), size=14, color=NAVY, bold=True, align=PP_ALIGN.RIGHT)
        top += Emu(520000)
    add_footer(ppt_slide, page_num, footer_text)


def render_comparison(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    columns = block_body(slide, kind="comparison-column")
    if len(columns) < 2:
        render_bullets(prs, slide, page_num, footer_text)
        return
    left_col, right_col = columns[:2]
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(ppt_slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(ppt_slide, slide["title"], Emu(365760), Emu(274320), Emu(7436815), Emu(457200), size=26, color=NAVY, bold=True)
    add_rect(ppt_slide, Emu(365760), Emu(1100000), Emu(3300000), Emu(420000), ACCENT_RED)
    add_rect(ppt_slide, Emu(4300000), Emu(1100000), Emu(3300000), Emu(420000), ACCENT_GREEN)
    add_text(ppt_slide, left_col["label"], Emu(365760), Emu(1180000), Emu(3300000), Emu(220000), size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    add_text(ppt_slide, right_col["label"], Emu(4300000), Emu(1180000), Emu(3300000), Emu(220000), size=15, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    left_points = [item.strip() for item in left_col["body"].split("|")]
    right_points = [item.strip() for item in right_col["body"].split("|")]
    top = Emu(1620000)
    for left_item, right_item in zip(left_points, right_points):
        add_rect(ppt_slide, Emu(365760), top, Emu(3300000), Emu(470000), RGBColor(0xFB, 0xF0, 0xF2), line=RGBColor(0xF5, 0xD7, 0xDF), radius=True)
        add_rect(ppt_slide, Emu(4300000), top, Emu(3300000), Emu(470000), RGBColor(0xEE, 0xFA, 0xF5), line=RGBColor(0xCC, 0xEE, 0xDE), radius=True)
        add_text(ppt_slide, f"×  {left_item}", Emu(520000), top + Emu(90000), Emu(2900000), Emu(260000), size=14, color=INK)
        add_text(ppt_slide, f"✓  {right_item}", Emu(4450000), top + Emu(90000), Emu(2900000), Emu(260000), size=14, color=INK)
        top += Emu(560000)
    add_footer(ppt_slide, page_num, footer_text)


def render_use_case_table(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    rows = block_body(slide, kind="table") or block_body(slide)
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(ppt_slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(ppt_slide, slide["title"], Emu(365760), Emu(274320), Emu(7600000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(ppt_slide, slide.get("objective", ""), Emu(365760), Emu(700000), Emu(8200000), Emu(220000), size=12, color=MUTED)
    top = Emu(1250000)
    for row in rows[:6]:
        add_rect(ppt_slide, Emu(365760), top, Emu(2200000), Emu(480000), TEAL_DARK, radius=True)
        add_text(ppt_slide, row["label"], Emu(520000), top + Emu(120000), Emu(1900000), Emu(220000), size=14, color=WHITE, bold=True)
        add_rect(ppt_slide, Emu(2700000), top, Emu(8100000), Emu(480000), RGBColor(0xF3, 0xF6, 0xF8), line=GRID)
        add_text(ppt_slide, row["body"], Emu(2880000), top + Emu(90000), Emu(7700000), Emu(280000), size=12, color=INK)
        top += Emu(620000)
    add_footer(ppt_slide, page_num, footer_text)


def render_snapshot_grid(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    cards = block_body(slide, kind="company-card")
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(ppt_slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(ppt_slide, slide["title"], Emu(365760), Emu(274320), Emu(7800000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(ppt_slide, slide.get("objective", ""), Emu(365760), Emu(700000), Emu(7600000), Emu(220000), size=12, color=MUTED)
    positions = [
        (Emu(365760), Emu(1300000)),
        (Emu(6200000), Emu(1300000)),
        (Emu(365760), Emu(3600000)),
        (Emu(6200000), Emu(3600000)),
    ]
    for card, (left, top) in zip(cards[:4], positions):
        add_rect(ppt_slide, left, top, Emu(5000000), Emu(1800000), CREAM, line=GRID, radius=True)
        add_text(ppt_slide, card["label"], left + Emu(170000), top + Emu(320000), Emu(3000000), Emu(240000), size=18, color=NAVY, bold=True)
        add_text(ppt_slide, card["body"], left + Emu(170000), top + Emu(760000), Emu(4500000), Emu(760000), size=12, color=INK)
    add_footer(ppt_slide, page_num, footer_text)


def render_case_study(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    cards = block_body(slide, kind="company-card")
    bullets = block_body(slide, kind="bullet-list") or block_body(slide, kind="callout")
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, TEAL)
    add_text(ppt_slide, slide["title"], Emu(365760), Emu(182880), Emu(9000000), Emu(502920), size=28, color=NAVY, bold=True)
    add_rect(ppt_slide, Emu(365760), Emu(658368), Emu(1828800), Emu(27432), NAVY)
    add_text(ppt_slide, "Evidence".upper(), Emu(365760), Emu(777240), Emu(2200000), Emu(200000), size=16, color=NAVY, bold=True)
    top = Emu(960120)
    for card in cards[:4]:
        add_text(ppt_slide, f"• {card['label']} — {card['body']}", Emu(365760), top, Emu(3520440), Emu(240000), size=13, color=INK)
        top += Emu(300000)
    add_text(ppt_slide, "Implication".upper(), Emu(4023360), Emu(777240), Emu(2200000), Emu(200000), size=16, color=NAVY, bold=True)
    top = Emu(960120)
    for bullet in bullets[:5]:
        add_text(ppt_slide, f"• {bullet['body']}", Emu(4023360), top, Emu(3520440), Emu(240000), size=13, color=INK)
        top += Emu(300000)
    add_footer(ppt_slide, page_num, footer_text)


def render_roadmap(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    items = block_body(slide, kind="roadmap-step") or block_body(slide)
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, CREAM)
    add_text(ppt_slide, slide["title"], Emu(731520), Emu(500000), Emu(8000000), Emu(500000), size=30, color=NAVY, bold=True)
    add_text(ppt_slide, slide.get("objective", ""), Emu(731520), Emu(980000), Emu(9000000), Emu(240000), size=13, color=MUTED)
    top = Emu(1650000)
    for item in items[:6]:
        label = item["label"]
        value = item["body"]
        if "|" in value:
            main, detail = [part.strip() for part in value.split("|", 1)]
        else:
            main, detail = value, ""
        add_rect(ppt_slide, Emu(900000), top, Emu(1200000), Emu(520000), TEAL, radius=True)
        add_text(ppt_slide, label, Emu(900000), top + Emu(140000), Emu(1200000), Emu(180000), size=13, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_rect(ppt_slide, Emu(2400000), top, Emu(3400000), Emu(520000), WHITE, line=GRID, radius=True)
        add_text(ppt_slide, main, Emu(2580000), top + Emu(120000), Emu(3000000), Emu(220000), size=16, color=NAVY, bold=True)
        add_rect(ppt_slide, Emu(6100000), top, Emu(4300000), Emu(520000), WHITE, line=GRID, radius=True)
        add_text(ppt_slide, detail, Emu(6280000), top + Emu(120000), Emu(3900000), Emu(220000), size=12, color=INK)
        top += Emu(760000)
    add_footer(ppt_slide, page_num, footer_text)


def render_bullets(prs: Presentation, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    ppt_slide = prs.slides.add_slide(slide_layout(prs))
    add_rect(ppt_slide, 0, 0, SLIDE_W, SLIDE_H, WHITE)
    add_rect(ppt_slide, 0, 0, SLIDE_W, Emu(70000), TEAL_DARK)
    add_text(ppt_slide, slide["title"], Emu(365760), Emu(274320), Emu(7600000), Emu(457200), size=26, color=NAVY, bold=True)
    add_text(ppt_slide, slide.get("objective", ""), Emu(365760), Emu(700000), Emu(8200000), Emu(260000), size=13, color=MUTED)
    top = Emu(1350000)
    for idx, block in enumerate(block_body(slide)[:7], start=1):
        add_rect(ppt_slide, Emu(500000), top, Emu(520000), Emu(520000), NAVY, radius=True)
        add_text(ppt_slide, f"{idx:02d}", Emu(500000), top + Emu(120000), Emu(520000), Emu(180000), size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
        add_rect(ppt_slide, Emu(1200000), top, Emu(9300000), Emu(620000), RGBColor(0xF3, 0xF6, 0xF8), line=GRID, radius=True)
        label = f"{block['label']}: " if block.get("label") else ""
        add_text(ppt_slide, label + block["body"], Emu(1450000), top + Emu(140000), Emu(8800000), Emu(340000), size=15, color=INK)
        top += Emu(760000)
    add_footer(ppt_slide, page_num, footer_text)


def render_slide(prs: Presentation, image_blob: bytes | None, slide: dict[str, Any], page_num: int, footer_text: str) -> None:
    slide_type = slide.get("slide_type", "custom")
    if slide_type == "hero":
        render_hero(prs, image_blob, slide, page_num, footer_text)
    elif slide_type in {"agenda"}:
        render_agenda(prs, slide, page_num, footer_text)
    elif slide_type in {"section-divider"}:
        render_section_divider(prs, slide, page_num, footer_text)
    elif slide_type in {"summary-cards", "recommendation"}:
        render_card_grid(prs, slide, page_num, footer_text)
    elif slide_type in {"kpi-grid"}:
        render_metric_grid(prs, slide, page_num, footer_text)
    elif slide_type in {"bar-chart", "distribution-chart"}:
        render_chart_like(prs, slide, page_num, footer_text)
    elif slide_type in {"comparison"}:
        render_comparison(prs, slide, page_num, footer_text)
    elif slide_type in {"use-case-table", "cluster-spotlight"}:
        render_use_case_table(prs, slide, page_num, footer_text)
    elif slide_type in {"snapshot-grid", "source-coverage"}:
        render_snapshot_grid(prs, slide, page_num, footer_text)
    elif slide_type in {"case-study"}:
        render_case_study(prs, slide, page_num, footer_text)
    elif slide_type in {"roadmap"}:
        render_roadmap(prs, slide, page_num, footer_text)
    else:
        render_bullets(prs, slide, page_num, footer_text)


def load_deck_plan(path: Path) -> dict[str, Any]:
    return json.loads(path.read_text(encoding="utf-8"))


def render_deck_to_file(deck_plan_path: Path, output_path: Path, template_path: Path = DEFAULT_TEMPLATE, footer_text: str | None = None) -> None:
    deck_plan = load_deck_plan(deck_plan_path)
    footer = footer_text or f"{deck_plan['report_slug']} | {deck_plan.get('audience', 'deck')}"
    image_blob = get_title_image(template_path)
    prs = build_presentation(template_path)
    for page_num, slide in enumerate(deck_plan["slides"], start=1):
        render_slide(prs, image_blob, slide, page_num, footer)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(output_path))


def main() -> None:
    parser = argparse.ArgumentParser(description="Render a deck-plan.json file to PPTX.")
    parser.add_argument("--deck-plan", type=Path, required=True)
    parser.add_argument("--output", type=Path, required=True)
    parser.add_argument("--template", type=Path, default=DEFAULT_TEMPLATE)
    parser.add_argument("--footer-text", default="")
    args = parser.parse_args()
    render_deck_to_file(args.deck_plan, args.output, args.template, args.footer_text or None)


if __name__ == "__main__":
    main()
